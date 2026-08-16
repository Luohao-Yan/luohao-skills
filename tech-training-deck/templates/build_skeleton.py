#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_skeleton.py — tech-training-deck 的 build 脚手架模板(从 profile.yaml 读品牌)。

这是 examples/deepseek-harness/build.py 的通用化版本:品牌色/字体/layout 从
profile.yaml 来(非硬编),含一套完整的 deck 节奏(封面/目录/章节页/内容页/
深色关键发现/红色结论/附录)+ 动画 + 讲者备注 + lint 门禁。

用法(复制到你的工作目录改内容):
    cp templates/build_skeleton.py my-topic/build.py
    # 改 TPL=你的.pptx、逐页内容、profile.yaml 的 semantic_contract
    python my-topic/build.py

需先跑 inspect_and_profile.py 生成 profile.yaml,并填好 semantic_contract。
依赖 slide-maker skill 的 deckkit/anim。
"""
import sys, os
# === 路径配置(改成你机器的) ===
SLIDE_MAKER = os.path.expanduser(r"~/.claude/skills/slide-maker/scripts")
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SLIDE_MAKER)
sys.path.insert(0, os.path.join(os.path.dirname(HERE), "scripts"))   # tech-training-deck/scripts
import deckkit as dk
from anim import Build
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from load_profile import load
from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes

# === 你的模板与输出 ===
TPL = r"<改为你自己的.pptx模板路径>"
OUT = os.path.join(HERE, "training-deck.pptx")

P = load(os.path.join(HERE, "profile.yaml"))
D = Deck(P)                       # 设好 dk.FONT/EAFONT + 语义色快捷
W_IN, H_IN = D.W, D.H
WHITE = RGBColor(0xFF,0xFF,0xFF)
CARDBG = RGBColor(0xF7,0xF7,0xFA)

def build():
    prs = dk.open_template(TPL)

    # -------- 1. 封面 (cover layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("cover")])
    try:
        t = s.placeholders[0]; t.text = "<主标题>"
        for p in t.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(40); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=dk.EAFONT
                dk._apply_ea(r, dk.EAFONT)
    except Exception: pass
    for pidx, txt in [(10,"<副标题:主题·关键词>"),(11,"<单位·日期>")]:
        try:
            ph = s.placeholders[pidx]; ph.text = txt
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = dk.EAFONT; dk._apply_ea(r, dk.EAFONT)
        except Exception: pass
    notes(s, "<开场:一句话讲清今天要讲什么、为什么对听众重要。>")

    # -------- 2. 目录 (content layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "目录  Catalog", D.anchor)
    col = dk.columns(2, slide=s, top=1.4, bottom=0.7, margin=0.62, gap=0.5)
    for label, items, accent in [("01 <第一部分>", ["<子项1>","<子项2>","<子项3>"], D.anchor),
                                   ("02 <第二部分>", ["<子项1>","<子项2>","<子项3>"], D.comparator)]:
        x,y,w,h = col[0] if label.startswith("01") else col[1]
        card(s, x, y, w, h, fill=CARDBG, line=accent, line_w=1.2)
        dk.text(s, x+0.3, y+0.22, w-0.6, 0.4, [[(label[:2], 22, accent, True, False, dk.FONT)]], wrap=False)
        dk.text(s, x+0.3, y+0.62, w-0.6, 0.4, [[(label[3:], 18, RGBColor(0x2A,0x2A,0x33), True, False, dk.EAFONT)]], wrap=False)
        yy = y+1.15
        for it in items:
            dk.text(s, x+0.4, yy, w-0.8, 0.32, [[("#  "+it, 13.5, RGBColor(0x2A,0x2A,0x33), False, False, dk.EAFONT)]])
            yy += 0.38
    notes(s, "<分几部分讲,为什么这个顺序。>")

    # -------- 3. 章节页 (chapter layout, 自动加深色衬底修对比) --------
    chap(prs, D, "chapter", "01", "<第一部分标题>", sub="<关键词列表>")
    notes(s, "<进入第一部分,为什么先讲这个。>")

    # -------- 4. 内容页范例 (content layout, 4-card) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<断言式标题:结论即标题>", D.anchor)
    cols = dk.columns(4, slide=s, top=1.35, bottom=0.95, margin=0.5, gap=0.28)
    data = [("<名1>","<归属1>","<形态1>","<要点1>", D.anchor),
            ("<名2>","<归属2>","<形态2>","<要点2>", D.comparator)]
    for (name,owner,form,feat,col),c in zip(data*2, cols):
        x,y,w,h = c
        card(s, x, y, w, h, fill=WHITE, line=col, line_w=1.6, r=0.1)
        dk.box(s, x, y, w, 0.52, fill=col, round=True, corners='top', r=0.1)
        dk.text(s, x+0.16, y+0.05, w-0.32, 0.42, [[(name, 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.18, y+0.64, w-0.36, 0.32, [[(owner, 13, col, True, False, dk.EAFONT)]], wrap=False)
        dk.text(s, x+0.18, y+1.0, w-0.36, 0.32, [[(form, 12.5, D.neutral, False, False, dk.EAFONT)]])
        dk.box(s, x+0.18, y+1.38, w-0.36, 0.012, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+0.18, y+1.5, w-0.36, 0.7, [[(feat, 13, RGBColor(0x2A,0x2A,0x33), False, False, dk.EAFONT)]], line_spacing=1.2)
    dk.bottom_callout(s, 0.5, W_IN-1.0, "<关键提示>", "<一句话纠偏或强调>", label_c=D.anchor, body_c=RGBColor(0x20,0x26,0x30))
    notes(s, "<这张表的讲解:逐卡点出归属/形态/要点,最后纠偏。>")

    # -------- 5. 深色关键发现页 (dark layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("dark")])
    set_title(s, "<关键发现:断言>", WHITE)
    cols = dk.columns(2, slide=s, top=1.35, bottom=0.75, margin=0.5, gap=0.5)
    lx = cols[0]; x,y,w,h = lx
    card(s, x, y, w, h, fill=RGBColor(0x2A,0x14,0x18), line=D.anchor, line_w=1.5)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.5, [[("本机实测", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.3, y+0.6, w-0.6, 1.2, [[("<发现内容>,",16,WHITE,True,False,dk.EAFONT)]], line_spacing=1.3)
    rx = cols[1]; x,y,w,h = rx
    card(s, x, y, w, h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=1.2)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.4, [[("意味着什么", 13, D.emphasis, True, False, dk.EAFONT)]], wrap=False)
    notes(s, "<这是最该让听众记住的发现,深色页做强调。>")

    # -------- 6. 结论 (red_conclusion layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("red_conclusion")])
    set_title(s, "结论", WHITE)
    dk.text(s, 1.0, 2.2, 11.3, 2.0, [
        [("<主题> 既是", 24, WHITE, True, False, dk.EAFONT),("技术借鉴对象", 26, D.emphasis, True, False, dk.EAFONT),
         (",也是", 24, WHITE, True, False, dk.EAFONT),("战略变量", 26, D.emphasis, True, False, dk.EAFONT),("。", 24, WHITE, True, False, dk.EAFONT)],
        [("", 8, WHITE, False, False)],
        [("短期:<...>", 16, WHITE, False, False, dk.EAFONT)],
        [("", 6, WHITE, False, False)],
        [("中期:<...>", 16, WHITE, False, False, dk.EAFONT)],
        [("", 6, WHITE, False, False)],
        [("长期:<...>", 16, WHITE, False, False, dk.EAFONT)],
    ], line_spacing=1.2)
    notes(s, "<收尾:一句话点睛 + 短中长期行动。>")

    # -------- 7. 附录 (content layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "附录·证据出处", D.anchor)
    notes(s, "<证据出处:按来源分组列 file_path:line。>")

    dk.lint_layout(prs, strict=True)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
