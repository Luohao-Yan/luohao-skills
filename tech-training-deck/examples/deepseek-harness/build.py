# -*- coding: utf-8 -*-
"""build.py — DeepSeek Harness 培训 deck(示例,品牌色从 profile.yaml 读,非硬编)

tech-training-deck skill 的 Stage 3 产物示例。复用 slide-maker 的 deckkit/anim,
品牌色/字体/layout 从 profile.yaml 通过 deck_helpers.Deck 加载。

跑通需:① 装好 slide-maker skill;② 把 TPL 改成你的 .pptx;③ profile.yaml 填好 semantic_contract。
"""
import sys, os
HERE = os.path.dirname(os.path.abspath(__file__))   # examples/deepseek-harness/
SKILL_ROOT = os.path.dirname(os.path.dirname(HERE))  # tech-training-deck/
sys.path.insert(0, os.path.join(SKILL_ROOT, "scripts"))   # 本 skill 的 scripts(含 slide_maker_path/load_profile/deck_helpers)
from slide_maker_path import find_slide_maker
sys.path.insert(0, find_slide_maker())   # slide-maker 的 scripts(deckkit/anim,自动探测)
import deckkit as dk
from anim import Build
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from load_profile import load
from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes, bottom_callout_at

TPL = r"<改为你自己的.pptx模板路径>"   # 例:r"C:/Users/.../某企业模板.pptx"
OUT = os.path.join(HERE, "training-deck.pptx")

P = load(os.path.join(HERE, "profile.yaml"))
D = Deck(P)
W_IN, H_IN = D.W, D.H
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2A,0x2A,0x33)
INK = RGBColor(0x2A,0x2A,0x33)
MUTE = RGBColor(0x6B,0x6B,0x6B)
CARDBG = RGBColor(0xF7,0xF7,0xFA)
DEEP = RGBColor(0xC6,0x00,0x00)

def build():
    prs = dk.open_template(TPL)

    # 1 封面
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("cover")])
    try:
        t = s.placeholders[0]; t.text = "DeepSeek Harness 能力培训"
        for p in t.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(40); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=dk.EAFONT
                dk._apply_ea(r, dk.EAFONT)
    except Exception: pass
    for pidx, txt in [(10,"给公司领导 · 一切皆插件 · 模型底座之争"),(11,"2026/08/16")]:
        try:
            ph = s.placeholders[pidx]; ph.text = txt
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = dk.EAFONT; dk._apply_ea(r, dk.EAFONT)
        except Exception: pass
    notes(s, "开场:今天用十几分钟,讲清 DeepSeek 开源的 Harness 是什么、为什么重要。")

    # 2 目录
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "目录  Catalog", D.anchor)
    col = dk.columns(2, slide=s, top=1.4, bottom=0.7, margin=0.62, gap=0.5)
    for label, items, accent, ci in [("01 认清四个产品",["四个产品到底是谁的","WorkBuddy/灵犀/Comate 区别","关键发现:Comate 内嵌 Pi"], D.anchor, 0),
                                       ("02 看懂 Harness 与战略",["一切皆插件怎么实现","DeepSeek 如何锁模型底座","dsh vs Pi 对比 · 给我们的建议"], D.comparator, 1)]:
        x,y,w,h = col[ci]
        card(s, x, y, w, h, fill=CARDBG, line=accent, line_w=1.2)
        dk.text(s, x+0.3, y+0.22, w-0.6, 0.4, [[(label[:2], 22, accent, True, False, dk.FONT)]], wrap=False)
        dk.text(s, x+0.3, y+0.62, w-0.6, 0.4, [[(label[3:], 18, DARK, True, False, dk.EAFONT)]], wrap=False)
        yy = y+1.15
        for it in items:
            dk.text(s, x+0.4, yy, w-0.8, 0.32, [[("#  "+it, 13.5, INK, False, False, dk.EAFONT)]]); yy += 0.38
    notes(s, "分两部分:先厘清易混产品,再看架构和战略,最后给建议。")

    # 3 章节页
    chap(prs, D, "chapter", "01", "先厘清:这些产品到底是什么", sub="WorkBuddy·灵犀·Comate·dsh 各是谁的")
    notes(s, "市面上几个产品常被混为一谈,先把归属讲清楚。")

    # 4 四产品归属表
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "先纠偏:四个产品,四个归属", D.anchor)
    cols = dk.columns(4, slide=s, top=1.35, bottom=1.15, margin=0.5, gap=0.28)
    data = [("DeepSeek Harness","DeepSeek 官方","开源 Agent 框架","编程底座\n不含办公",D.anchor),
            ("WorkBuddy","腾讯·非DeepSeek","桌面工作台","编程强\n办公靠技能",D.comparator),
            ("WPS 灵犀","金山 WPS","AI 办公智能体","编程+办公\n都在壳里",DEEP),
            ("WPS Comate","金山 WPS","桌面 Agent","内嵌 Pi\n编程+JSAPI",D.neutral)]
    for (name,owner,form,feat,col),c in zip(data, cols):
        x,y,w,h = c
        card(s, x, y, w, h, fill=WHITE, line=col, line_w=1.6, r=0.1)
        dk.box(s, x, y, w, 0.52, fill=col, round=True, corners='top', r=0.1)
        dk.text(s, x+0.16, y+0.05, w-0.32, 0.42, [[(name, 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.18, y+0.64, w-0.36, 0.32, [[(owner, 13, col, True, False, dk.EAFONT)]], wrap=False)
        dk.text(s, x+0.18, y+1.0, w-0.36, 0.32, [[(form, 12.5, D.neutral, False, False, dk.EAFONT)]])
        dk.box(s, x+0.18, y+1.38, w-0.36, 0.012, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+0.18, y+1.5, w-0.36, 0.7, [[(feat, 13, INK, False, False, dk.EAFONT)]], line_spacing=1.2)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "关键纠偏", "WorkBuddy 是腾讯的、非 DeepSeek;灵犀(千万办公)才含 WPS 办公功能;Comate 内嵌第三方开源 Pi。", label_c=D.anchor, body_c=DARK)
    notes(s, "WorkBuddy 不是 DeepSeek 的,是腾讯;灵犀才真正含 WPS 办公;Comate 内嵌第三方 Pi。")

    # 5 关键发现(深色页)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("dark")]); set_title(s, "关键发现:WPS Comate 内嵌的就是 Pi", WHITE)
    cols = dk.columns(2, slide=s, top=1.35, bottom=1.3, margin=0.5, gap=0.5)
    x,y,w,h = cols[0]
    card(s, x, y, w, h, fill=RGBColor(0x2A,0x14,0x18), line=D.anchor, line_w=1.5)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.5, [[("本机实测", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.3, y+0.6, w-0.6, 1.2, [[("金山自家的 WPS Comate,",16,WHITE,True,False,dk.EAFONT),("核心就是 Pi",16,D.emphasis,True,False,dk.EAFONT),(" v0.79.3",15,WHITE,False,False,dk.EAFONT)]], line_spacing=1.3)
    x,y,w,h = cols[1]
    card(s, x, y, w, h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=1.2)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.4, [[("意味着什么", 13, D.emphasis, True, False, dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["Pi 轻量路线已被产品化落地","选型与产品方向一致、风险低","与 dsh 同源 pi-ai,可平滑评估"]):
        dk.text(s, x+0.3, y+0.7+i*0.56, w-0.6, 0.5, [[("▸  ",13,D.anchor,True,False,dk.FONT),(it,13,WHITE,False,False,dk.EAFONT)]], line_spacing=1.2)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "三大 harness 互不内嵌", "WorkBuddy 用 codebuddy、Comate 用 Pi、DeepSeek 用 dsh——各自独立。", label_c=D.emphasis, body_c=WHITE, fill=RGBColor(0x2A,0x14,0x18))
    notes(s, "金山自家 Comate 的核心运行时就是第三方 Pi;我们也在用 Pi。三大 harness 互不内嵌。")

    # 6 章节页
    chap(prs, D, "chapter", "02", 'DeepSeek Harness 如何"一切皆插件"', sub="Cordis·能力 seam·插件树·创造模式")
    notes(s, "进入架构,目标让领导看懂插件化体现在哪。")

    # 7 定位+四模式
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "定位:Agent = 模型 + Harness", D.anchor)
    dk.text(s, 0.7, 1.5, 12, 1.2, [[("Agent",40,D.anchor,True,False,dk.FONT),(" = ",32,DARK,True,False,dk.FONT),("模型",34,D.neutral,True,False,dk.EAFONT),(" + ",30,DARK,True,False,dk.FONT),("Harness",40,D.comparator,True,False,dk.FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, 0.7, 2.7, 12, 0.5, [[("模型负责『想』,Harness 负责理解环境、调工具、组织多步任务", 14, MUTE, False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
    modes = [("标准","默认全套工具,日常开发"),("PTC","生成代码组合多步调用"),("极简","只留 Shell,跑基准"),("创造","agent 改自己的插件")]
    cols = dk.columns(4, slide=s, top=3.6, bottom=1.3, margin=0.5, gap=0.3)
    b8 = Build(s)
    for i,((n,d),c) in enumerate(zip(modes, cols)):
        x,y,w,h = c
        with b8.step():
            card(s, x, y, w, h, fill=CARDBG, line=D.neutral, line_w=1.0, r=0.1)
            num_circle(s, x+w/2-0.28, y+0.18, 0.56, i+1, D.anchor, D.comparator, size=16)
            dk.text(s, x+0.12, y+0.82, w-0.24, 0.35, [[(n, 15, D.anchor, True, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
            dk.text(s, x+0.12, y+1.18, w-0.24, 0.6, [[(d, 14, INK, False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    b8.apply(effect="fade")
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "让大模型『能动手干活』的运行时底座,对标 Claude Code / Codex,开源。", label_c=D.anchor, body_c=DARK)
    notes(s, "Agent=模型+Harness。四模式:标准/PTC/极简/创造。创造模式 agent 能改自己的插件。")


    # 7b 整体架构全景(五层)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "整体架构：五层插件，自下而上组合", D.anchor)
    layers = [("Cordis 内核 (vendor)","插件元框架：ctx 容器 · inject · 事件 · 注册可逆",D.neutral,RGBColor(0xE9,0xEE,0xF5)),("核心服务 seams","ctx.llm · ctx.tools · ctx.agents · ctx.shell · ctx.fs · ctx.sessions · ctx.subagents",D.neutral,RGBColor(0xF2,0xF5,0xF9)),("能力插件包 (packages/*)","llm-deepseek/llm-pi-ai · tool-* · shell-* · fs · lsp · subagent-* · web · skill · workflow",D.anchor,RGBColor(0xFF,0xF4,0xF6)),("组合包 / 组合点","dsh-base bundle · preset (cordis.yml) · profile (用户 patch)",D.anchor,RGBColor(0xFF,0xF1,0xF3)),("apps 组装点","apps/cli (dsh 命令) · apps/web · ACP · JSON-RPC",RGBColor(0xC6,0,0),RGBColor(0xFF,0xEE,0xF0))]
    lx, lw = 0.7, 11.9; ly0 = 1.45; lh = 0.78; gap = 0.18
    for i,(t,d,col,fill) in enumerate(layers):
        y = ly0 + i*(lh+gap)
        card(s, lx, y, lw, lh, fill=fill, line=col, line_w=1.6 if i>=2 else 1.0, r=0.08)
        dk.box(s, lx, y, 0.1, lh, fill=col, round=True, corners='left', r=0.08)
        dk.text(s, lx+0.25, y+0.05, 3.6, 0.34, [[(t,14,col,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, lx+3.95, y+0.05, lw-4.15, lh-0.1, [[(d,12,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        if i < len(layers)-1: dk.arrow(s, lx+lw/2-0.12, y+lh+0.0, 0.24, gap, color=RGBColor(0xB0,0xB8,0xC4), direction='down')
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "host/client 双面", "core 既是被组合的插件集，也暴露 host/client 给外部进程（ACP/JSON-RPC）——同一内核多种形态。", label_c=D.anchor, body_c=RGBColor(0x20,0x26,0x30))
    notes(s, "dsh 五层架构:Cordis内核→核心seams→能力插件包→组合包→apps。组合↑依赖↓。host/client双面。证据:packages/README.zh.md、bundle/README.zh.md、pnpm-workspace.yaml。")

    # 8 Cordis 五概念 (架构示意图:ctx 容器 + 4 插件挂载)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "Cordis 五概念：插件挂到 ctx，按 key 找、可撤销", D.anchor)
    ctx_x, ctx_y, ctx_w, ctx_h = 0.7, 1.55, 3.4, 4.9
    dk.box(s, ctx_x, ctx_y, ctx_w, ctx_h, fill=RGBColor(0xF3,0xF5,0xF8), line=D.neutral, line_w=2.0, round=True, r=0.1)
    dk.box(s, ctx_x+0.2, ctx_y+0.2, ctx_w-0.4, 0.62, fill=D.neutral, round=True, r=0.08)
    dk.text(s, ctx_x+0.2, ctx_y+0.2, ctx_w-0.4, 0.62, [[("插件 = Service 对象",13.5,WHITE,True,False,dk.EAFONT),("  (apply ctx)",11,RGBColor(0xCC,0xDD,0xEE),False,False,dk.FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, ctx_x+0.2, ctx_y+0.92, ctx_w-0.4, 0.3, [[("ctx 服务容器 · 按稳定 key 持有",11,RGBColor(0x6B,0x6B,0x6B),False,False,dk.EAFONT)]], align=PP_ALIGN.CENTER)
    pairs = [("ctx.tools/llm","按 key 找","按稳定 key 找服务,不 import 实现"),("inject:[]","inject 依赖","等依赖就绪才启动,顺序由依赖决定"),("emit/waterfall/…","事件四模式","观察·改写·扇出·按序"),("ctx.effect/on","注册可逆","装的副作用卸载时自动撤销")]
    slot_x = ctx_x+0.3; slot_w = ctx_w-0.6
    slot_ys = [ctx_y+1.35 + i*0.82 for i in range(4)]
    plug_x = 5.6; plug_w = 7.1; plug_ys = [1.75 + i*1.02 for i in range(4)]
    for i,(ky,n,d) in enumerate(pairs):
        sy = slot_ys[i]
        dk.box(s, slot_x, sy, slot_w, 0.6, fill=WHITE, line=RGBColor(0xC8,0xD0,0xDC), line_w=1.0, round=True, r=0.08)
        dk.text(s, slot_x+0.15, sy+0.03, slot_w-0.3, 0.55, [[(ky,12.5,D.anchor,True,False,dk.FONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        py = plug_ys[i]
        card(s, plug_x, py, plug_w, 0.86, fill=WHITE, line=D.anchor, line_w=1.4, r=0.08)
        dk.box(s, plug_x, py, 0.12, 0.86, fill=D.anchor, round=True, corners='left', r=0.08)
        dk.text(s, plug_x+0.25, py+0.1, 2.4, 0.34, [[(n,14,D.anchor,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, plug_x+0.25, py+0.46, 2.4, 0.3, [[(ky,11,D.neutral,False,False,dk.FONT)]], wrap=False)
        dk.box(s, plug_x+2.75, py+0.14, 0.012, 0.6, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, plug_x+2.9, py+0.08, plug_w-3.05, 0.72, [[(d,13,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        dk.connector(s, (slot_x+slot_w, sy+0.3), (plug_x, py+0.43), style="solid", color=RGBColor(0xB0,0xB8,0xC4), width=1.5, arrow=True)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "没有特权内核——连 agent 循环本身都是 ctx 上的一个插件,随时可换可拆。", label_c=D.anchor, body_c=RGBColor(0x20,0x26,0x30))
    notes(s, "Cordis 五概念:插件=Service对象/按key找/inject依赖/事件四模式/注册可逆。架构图:ctx容器持key,插件挂到key上。没有特权内核。")


    # 8b Agent loop 执行流程
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "Agent loop：一轮 turn 怎么跑（ReactLoopAgent）", D.anchor)
    stages = [("1 构造请求","buildRequest","上下文 + compaction 压缩 + 工具 schema 注入",D.neutral),("2 调 LLM","ctx.llm","经 seam 调模型（deepseek/pi-ai/replay 可换）→ 流式输出 + tool calls",D.anchor),("3 解析工具调用","ctx.tools","tool calls 落到 ctx.tools 上对应工具",D.neutral),("4 工具执行管线","审批→沙箱→执行","interaction 审批 → sandbox → 执行 → 结果喂回上下文",D.anchor),("5 判停","agent/turn-stopping","goal 达成 / 用户中断 / max turns / 不再请求工具",RGBColor(0xC6,0,0))]
    sx, sw = 0.7, 11.9; sy0 = 1.5; sh = 0.72; gap = 0.2
    for i,(t,fn,d,col) in enumerate(stages):
        y = sy0 + i*(sh+gap)
        card(s, sx, y, sw, sh, fill=WHITE, line=col, line_w=1.4, r=0.08)
        dk.box(s, sx, y, 0.1, sh, fill=col, round=True, corners='left', r=0.08)
        dk.text(s, sx+0.25, y+0.04, 2.5, 0.34, [[(t,14,col,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, sx+0.25, y+0.38, 2.5, 0.3, [[(fn,11,D.neutral,False,False,dk.FONT)]], wrap=False)
        dk.box(s, sx+2.85, y+0.1, 0.012, sh-0.2, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, sx+3.0, y+0.04, sw-3.2, sh-0.08, [[(d,12.5,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        if i < len(stages)-1: dk.arrow(s, sx+sw/2-0.12, y+sh+0.0, 0.24, gap, color=RGBColor(0xB0,0xB8,0xC4), direction='down')
    dk.text(s, sx, sy0+len(stages)*(sh+gap)-0.04, sw, 0.3, [[("未停则 continue 下一 turn · 主循环对所有模式一致（模式只改工具集+prompt）",11,D.neutral,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER)
    notes(s, "Agent loop:ReactLoopAgent逐turn。五阶段:buildRequest→ctx.llm→解析tool calls→执行管线(审批→沙箱→执行→喂回)→turn-stopping判停。模式不改loop只改工具集+prompt。证据:packages/core/agent/、compaction、interaction。")

    # 9 signature move: ctx.llm 插槽图
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "能力 seam:换一个 Provider,换一个世界", D.anchor)
    slot_x, slot_y, slot_w, slot_h = 0.62, 1.4, 3.3, 2.6
    card(s, slot_x, slot_y, slot_w, slot_h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=2.0, r=0.1)
    dk.text(s, slot_x+0.2, slot_y+0.15, slot_w-0.4, 0.4, [[("ctx.llm 插槽", 14, D.emphasis, True, False, dk.FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.box(s, slot_x+0.35, slot_y+1.1, slot_w-0.7, 1.2, fill=RGBColor(0x0D,0x12,0x1A), line=D.anchor, line_w=1.5, round=True, r=0.08)
    dk.text(s, slot_x+0.35, slot_y+1.1, slot_w-0.7, 1.2, [[("Provider\n插这里",14,RGBColor(0xFF,0xAA,0xAA),False,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    dk.arrow(s, slot_x+slot_w+0.05, slot_y+slot_h/2-0.25, 0.7, 0.5, color=D.anchor, direction='right')
    provs = [("llm-deepseek","DeepSeek 原生",D.anchor,True),("llm-pi-ai","经 pi-ai 接任意",D.comparator,False),("llm-replay","回放测试",MUTE,False)]
    px = slot_x+slot_w+0.95; pw = (W_IN-0.62-px-0.3)/3 - 0.2
    for i,(nm,d,col,hi) in enumerate(provs):
        x = px + i*(pw+0.25); y = slot_y+0.1
        card(s, x, y, pw, 2.6, fill=RGBColor(0xFF,0xF4,0xF6) if hi else WHITE, line=col, line_w=2.2 if hi else 1.2, r=0.1)
        if hi:
            dk.box(s, x, y, pw, 0.4, fill=col, round=True, corners='top', r=0.1)
            dk.text(s, x+0.1, y+0.02, pw-0.2, 0.36, [[("★ 默认",11,WHITE,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.12, y+0.5, pw-0.24, 0.4, [[(nm,14,col,True,False,dk.FONT)]], align=PP_ALIGN.CENTER, wrap=False)
        dk.text(s, x+0.12, y+0.95, pw-0.24, 0.5, [[(d,13,INK,False,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.15)
        dk.text(s, x+0.12, y+1.62, pw-0.24, 0.8, [[("换这个\n= 换模型底座",14,D.anchor,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "seam 三角色", "Service Definition·Service Provider(并列可换)·Consumer——三者一并设计才是一个完整能力。", label_c=D.anchor, body_c=DARK)
    notes(s, "ctx.llm 是插槽,三个 Provider 可换,DeepSeek 默认。换 Provider=换模型底座。")


    # 9b extensions 生态与插件分发
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "生态：插件即 npm/git 包，dsh plugin add 安装", D.anchor)
    lx, lw = 0.7, 6.7
    card(s, lx, 1.5, lw, 4.9, fill=WHITE, line=D.neutral, line_w=1.2, r=0.1)
    dk.box(s, lx, 1.5, lw, 0.5, fill=D.neutral, round=True, corners='top', r=0.1)
    dk.text(s, lx+0.2, 1.55, lw-0.4, 0.4, [[("分发机制",14,WHITE,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    items = [("bundle / patch-layer","Cordis 配置+挂载代码的分发格式，可被上层 patch"),("dsh plugin add","dsh plugin --profile <name> add <package-or-git-spec>"),("dsh-plugin GitHub topic","开发者给插件仓库打此标签被发现（去中心化）"),("官方 extensions 包","packages/extensions/ 扩展包都是 Cordis 插件，一等公民")]
    for i,(k,d) in enumerate(items):
        yy = 2.2 + i*1.0
        dk.text(s, lx+0.25, yy, lw-0.5, 0.32, [[(k,13,D.anchor,True,False,dk.EAFONT)]], wrap=False)
        dk.text(s, lx+0.25, yy+0.36, lw-0.5, 0.55, [[(d,11.5,INK,False,False,dk.EAFONT)]], line_spacing=1.15)
    rx, rw = 7.7, 5.1
    card(s, rx, 1.5, rw, 2.4, fill=RGBColor(0xFF,0xF4,0xF6), line=RGBColor(0xC6,0,0), line_w=1.4, r=0.1)
    dk.box(s, rx, 1.5, rw, 0.5, fill=RGBColor(0xC6,0,0), round=True, corners='top', r=0.1)
    dk.text(s, rx+0.2, 1.55, rw-0.4, 0.4, [[("时效提醒",14,WHITE,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, rx+0.25, 2.1, rw-0.5, 1.7, [[("旧的 .dsh-plugin 仓库插件市场",12.5,RGBColor(0xC6,0,0),True,False,dk.EAFONT),(" 已于 2026-08-09 移除。以 ",12.5,INK,False,False,dk.EAFONT),("dsh-plugin GitHub topic + dsh plugin add",12.5,D.anchor,True,False,dk.FONT),(" 为准。",12.5,INK,False,False,dk.EAFONT)]], line_spacing=1.3)
    card(s, rx, 4.05, rw, 2.35, fill=WHITE, line=D.anchor, line_w=1.2, r=0.1)
    dk.box(s, rx, 4.05, rw, 0.5, fill=D.anchor, round=True, corners='top', r=0.1)
    dk.text(s, rx+0.2, 4.1, rw-0.4, 0.4, [[("运行时自修改 = 生态一环",13,WHITE,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, rx+0.25, 4.65, rw-0.5, 1.7, [[("创造模式 agent 经 cordis-host-runner 临时挂载/卸载插件——",11.5,INK,False,False,dk.EAFONT),("等于 agent 给自己装临时插件",11.5,D.anchor,True,False,dk.EAFONT),("，同属一切皆插件生态。",11.5,INK,False,False,dk.EAFONT)]], line_spacing=1.3)
    notes(s, "dsh 生态去中心化:插件即npm/git包,dsh plugin add安装,dsh-plugin topic发现。旧.dsh-plugin市场2026-08-09移除。运行时自修改也属生态。证据:packages/bundle/、packages/extensions/、dsh plugin add。")

    # 10 章节页
    chap(prs, D, "chapter", "03", "DeepSeek 如何用 Harness 锁定模型底座", sub="用『动手层』标准化,反向锁定『模型层』")
    notes(s, "进入战略,今天最该记住的部分。")

    # 11 战略四步
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "战略:用『动手层』标准化,反向锁定『模型层』", D.anchor)
    steps = [("默认带 DeepSeek","模型在 dsh 是可换插件,但默认体验/生态偏 DeepSeek"),("先选 harness 再选模型","团队建技术栈后迁移成本陡增;换模型只是一个插件"),("数据/反馈回流","更多 Agent 跑在 DeepSeek 上,真实数据反哺迭代"),("生态飞轮","Cordis + dsh-plugin 话题;24h 已 288 个插件仓")]
    rows = dk.rows(4, slide=s, top=1.4, bottom=1.25, gap=0.22)
    b13 = Build(s)
    for i,((n,d),c) in enumerate(zip(steps, rows)):
        x,y,w,h = c
        with b13.step():
            num_circle(s, x+0.1, y+h/2-0.27, 0.54, i+1, D.anchor, D.comparator, size=18)
            card(s, x+0.8, y, w-0.8, h, fill=CARDBG, line=RGBColor(0xE8,0xE8,0xEE), line_w=1.0, r=0.08)
            dk.text(s, x+1.0, y+0.12, w-1.2, 0.35, [[(n,15,D.anchor,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            dk.text(s, x+1.0, y+0.12, w-1.2, h-0.24, [[("　",1,WHITE,False,False),(d,14,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    b13.apply(effect="fade")
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "谁定义 harness 标准,谁就定义模型被怎么用、用谁的——DeepSeek 同时开源模型与 harness,用动手层标准化反向锁定模型层。", label_c=D.anchor, body_c=DARK)
    notes(s, "四步:默认带DeepSeek/先选harness/数据回流/生态飞轮。谁定义harness谁定义模型分发。")

    # 12 dsh vs Pi 对比
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "dsh 与 Pi:同源 pi-ai,分道于架构", D.anchor)
    headers = ["维度","DeepSeek Harness (dsh)","Pi"]
    rows_data = [["定位","官方开源·有 Web UI+CLI","独立开发者·纯终端 TUI(已入 Comate)"],["架构","一切皆插件·Cordis DI 容器","最小核心·注册式扩展"],["模型","默认 DeepSeek·多 provider","不锁定·30+ provider"],["能力包","官方维护几十个(一等公民)","核心 5 包,余靠扩展"],["办公","无","无(Comate 经 JSAPI 补)"],["部署","npx dsh web 一键起","npm i -g pi-coding-agent"],["底层","把 pi-ai 当可插拔后端","把 pi-ai 当内置底座"]]
    col_w = [1.7, 5.75, 5.25]
    dk.table(s, 0.62, 1.4, sum(col_w), [headers]+rows_data, col_w=col_w, header=True, size=13, row_h=0.56, head_c=D.anchor, body_c=INK, hi_fill=RGBColor(0xFF,0xF4,0xF6), hi_c=D.anchor)
    dk.text(s, 0.62, H_IN-0.55, 12.1, 0.45, [[("关键同源:",12.5,D.anchor,True,False,dk.EAFONT),("两者都用 ",12.5,INK,False,False,dk.EAFONT),("@earendil-works/pi-ai",12.5,D.neutral,True,False,dk.FONT),(" 作 LLM 抽象层——用 Pi 调通公司模型的经验可复用到评估 dsh。",12.5,INK,False,False,dk.EAFONT)]], line_spacing=1.15)
    notes(s, "对比:Pi 非金山自研但已被 Comate 内嵌。七维对比。同源 pi-ai,迁移成本低。")

    # 13 建议
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "五条建议,应对 harness 之争与模型趋同", D.anchor)
    advs = [("跟踪分发权","把 harness=分发权列为新战略变量","政企选型可能先选框架再选模型"),("借鉴 seam 化","研究 dsh seam 对架构的借鉴","适配私有化换模型/换部署/换沙箱"),("Pi 路线坚定","Pi 已被产品化验证、与 dsh 同源","短期提效,dsh 作平台化备选"),("勿上生产","dsh 是开发者预览、有破坏性变更","定位预研/POC,版本收敛再评估"),("厘清口径","对外统一产品归属口径","讲混会损害可信度")]
    rows = dk.rows(5, slide=s, top=1.4, bottom=0.75, gap=0.2)
    for i,((n,d,how),c) in enumerate(zip(advs, rows)):
        x,y,w,h = c
        num_circle(s, x, y+h/2-0.24, 0.48, i+1, D.anchor, D.comparator, size=15)
        card(s, x+0.65, y, w-0.65, h, fill=WHITE if i%2==0 else CARDBG, line=RGBColor(0xE8,0xE8,0xEE), line_w=1.0, r=0.06)
        dk.text(s, x+0.85, y+0.05, 3.1, h-0.1, [[(n,14,D.anchor,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.box(s, x+4.05, y+0.1, 0.012, h-0.2, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+4.25, y+0.05, w-4.45, h-0.1, [[(d,14,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    dk.text(s, 0.62, H_IN-0.42, 12, 0.32, [[("口径:",11,MUTE,False,False,dk.EAFONT),("WorkBuddy=腾讯·灵犀=金山·Comate=金山(内嵌Pi)·Pi=第三方·dsh=DeepSeek",11,D.neutral,False,False,dk.EAFONT)]])
    notes(s, "五条:跟踪分发权/借鉴seam/Pi坚定/勿上生产/厘清口径。")

    # 14 结论
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("red_conclusion")]); set_title(s, "结论", WHITE)
    dk.text(s, 1.0, 2.2, 11.3, 2.0, [[("dsh 既是",24,WHITE,True,False,dk.EAFONT),("技术借鉴对象",26,D.emphasis,True,False,dk.EAFONT),(",也是",24,WHITE,True,False,dk.EAFONT),("战略变量",26,D.emphasis,True,False,dk.EAFONT),("。",24,WHITE,True,False,dk.EAFONT)],[("",8,WHITE,False,False)],[("短期:Pi 路线已被验证落地、且本团队在用——可坚定。",16,WHITE,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("中期:dsh 作平台化底座候选持续观察,借鉴其 seam 化设计。",16,WHITE,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("长期:评估把自身方案架构『接缝化』,对冲模型层趋同风险。",16,WHITE,False,False,dk.EAFONT)]], line_spacing=1.2)
    notes(s, "收尾:dsh 既是借鉴对象也是战略变量。短期Pi坚定,中期dsh观察,长期接缝化。")

    # 15 附录
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "附录·证据出处", D.anchor)
    dk.text(s, 0.7, 1.5, 12, 5, [[("A. DeepSeek Harness(本地源码):",12.5,D.anchor,True,False,dk.EAFONT),(" README/architecture/cordis-primer · agent-presets · cordis-host-runner · llm-pi-ai",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("B. Pi(本机已装 0.84.2):",12.5,D.comparator,True,False,dk.EAFONT),(" @earendil-works/pi-coding-agent · 本机 models.json · 与 dsh 共用 pi-ai",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("C. WorkBuddy(本机实测):",12.5,D.anchor,True,False,dk.EAFONT),(" AppData/Local/Programs/WorkBuddy(腾讯,Electron)·codebuddy CLI·copilot.tencent.com",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("D. WPS Comate(本机实测):",12.5,D.anchor,True,False,dk.EAFONT),(" AppData/Local/WPS Comate(金山·内嵌Pi v0.79.3)·comate.wps.cn",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("E. WPS 灵犀(官方):",12.5,D.anchor,True,False,dk.EAFONT),(" lingxi.wps.cn · Python/Node/AirScript + WPS Office 集成",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("完整证据见同名 md 培训稿附录。",11,MUTE,False,False,dk.EAFONT)]], line_spacing=1.1)
    notes(s, "附录:所有结论的证据出处。完整证据链在培训稿 md 附录。")

    dk.lint_layout(prs, strict=True)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
