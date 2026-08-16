#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""new_deck.py — 给定 profile.yaml + 页数大纲,生成一个 build 脚手架脚本是 Stage 3 的起点:同事跑完 inspect_and_profile 拿到 profile.yaml 后,
用它生成一个可编辑的 build_<topic>.py(已配好:deckkit import、profile 加载、
deck_helpers、封面/目录/章节页/内容页/结论/附录的占位结构、lint 门禁、save)。
同事只需把每页的占位文字换成自己的内容。

用法:
    python new_deck.py --profile profile.yaml --topic "我的技术主题" --pages 8 \\
        --out build_my-topic.py
    python build_my-topic.py    # 生成 deck;需先有用户 .pptx 模板路径
"""
import argparse, os, textwrap

SKELETON = '''# -*- coding: utf-8 -*-
"""build_{slug}.py — {topic}(给领导培训 deck)
模板分支:从 profile.yaml 读品牌色/字体,不硬编。
Source of truth for the deck; re-run to rebuild identically.
依赖 slide-maker skill 的 deckkit/anim(import 路径见下)。
"""
import sys, os
HERE = os.path.dirname(os.path.abspath(__file__))
# === 改成你机器上 tech-training-deck skill 的 scripts 目录 ===
# (装了本 skill后,通常是 ~/.claude/skills/tech-training-deck/scripts 或 ~/.agents/skills/tech-training-deck/scripts)
SKILL_SCRIPTS = r"<改成 tech-training-deck skill 的 scripts 目录>"
sys.path.insert(0, SKILL_SCRIPTS)
from slide_maker_path import find_slide_maker
sys.path.insert(0, find_slide_maker())   # slide-maker 的 scripts(deckkit/anim,自动探测)
import deckkit as dk
from anim import Build
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from load_profile import load
from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes

# --- 路径:你的模板与输出 ---
TPL = r"<改为你自己的.pptx模板路径>"
OUT = os.path.join(HERE, "{slug}.pptx")

P = load(os.path.join(HERE, "profile.yaml"))
D = Deck(P)          # 设好 deckkit 全局字体 + 语义色快捷

def build():
    prs = dk.open_template(TPL)
    # TODO: 按 --pages 大纲逐页实现。每页节奏:
    #   s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    #   set_title(s, "断言式标题", D.anchor)
    #   cols = dk.columns(n, slide=s, top=1.4, bottom=0.95, margin=0.5, gap=0.3)
    #   for ...: card(...); dk.text(...); num_circle(...)
    #   notes(s, "讲者话术...")
    # 详见 examples/deepseek-harness/build.py 的完整范例。
{PAGE_STUBS}
    dk.lint_layout(prs, strict=True)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
'''

PAGE_STUB = '''
    # -------- {n}. {label} --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("{role}")])
    set_title(s, "{label}", D.anchor)
    notes(s, "{label} 的讲者话术:把详细内容放这里,幻灯片只留短语。")
'''

def slugify(t):
    import re
    return re.sub(r"[^a-z0-9_-]+", "-", t.lower()).strip("-") or "deck"

def main():
    ap = argparse.ArgumentParser(description="生成 build 脚手架")
    ap.add_argument("--profile", default="profile.yaml")
    ap.add_argument("--topic", required=True, help="主题名(中文OK)")
    ap.add_argument("--pages", type=int, default=8, help="页数(含封面/目录/附录)")
    ap.add_argument("--out", default=None)
    args = ap.parse_args()
    slug = slugify(args.topic)
    out = args.out or f"build_{slug}.py"
    # 简单页序:封面 + 目录 + N内容 + 结论 + 附录
    stubs = []
    stubs.append(PAGE_STUB.format(n=1, label="封面", role="cover"))
    stubs.append(PAGE_STUB.format(n=2, label="目录", role="content"))
    n_content = max(1, args.pages - 3)
    for i in range(n_content):
        stubs.append(PAGE_STUB.format(n=3+i, label=f"内容页{i+1}", role="content"))
    stubs.append(PAGE_STUB.format(n=3+n_content, label="结论", role="red_conclusion"))
    stubs.append(PAGE_STUB.format(n=4+n_content, label="附录·证据出处", role="content"))
    code = SKELETON.format(slug=slug, topic=args.topic, PAGE_STUBS="".join(stubs))
    open(out, "w", encoding="utf-8").write(code)
    print(f"-> {out}")
    print("next: 改 TPL 路径为你的 .pptx,逐页填内容,然后 python {out}".replace("{out}", out))

if __name__ == "__main__":
    main()
