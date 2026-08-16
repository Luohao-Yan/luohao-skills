#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""deck_helpers.py — 在模板上构建 deck 的通用积木,颜色/字体从 profile 来(非硬编)。

build 脚本:
    import sys, os
    sys.path.insert(0, find_slide_maker())          # deckkit/anim 在这(自动探测)
    sys.path.insert(0, "<tech-training-deck>/scripts")
    import deckkit as dk
    from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes

这些 helper 把"在用户模板上画一页"的重复动作封装好:填标题占位符并打 CJK 字体
tag、画模板语汇的渐变编号圆、画带半透明衬底的章节页(修深暖背景白字对比坑)、
画卡片、写讲者备注。颜色都从传入的 Profile 来,不硬编 hex。
"""
import os, sys
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import deckkit as dk

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD_FALLBACK = RGBColor(0xFF, 0xC0, 0x00)

class Deck:
    """持有 profile + deckkit 全局字体,build 脚本从这取颜色/字体/layout。"""
    def __init__(self, profile):
        self.P = profile
        latin, ea = profile.fonts()
        dk.FONT = latin; dk.EAFONT = ea; dk.DISPLAY = latin
        self.W, self.H = profile.canvas()
    # 语义色快捷(走 profile 的 semantic_contract)
    @property
    def anchor(self): return self.P.color("anchor_subject")
    @property
    def comparator(self): return self.P.color("comparator")
    @property
    def neutral(self): return self.P.color("neutral")
    @property
    def emphasis(self): return self.P.color("emphasis")

def set_title(slide, text, color, size=26, bold=True, font=None, ea=None):
    """填 idx=0 标题占位符 + 设字号/粗/色 + 打 CJK ea 字体 tag(kinsoku/渲染才正确)。
    每个内容页都要做;颜色由调用方从 profile 传(通常是 anchor 主色)。"""
    ph = slide.placeholders[0]
    ph.text = text
    latin, eaf = (font or dk.FONT), (ea or dk.EAFONT)
    for p in ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = latin
            dk._apply_ea(r, eaf)
    return ph

def num_circle(slide, cx, cy, d, num, color_a, color_b, tcolor=WHITE, size=16, font=None):
    """渐变编号圆(模板语汇):color_a→color_b 径向渐变 + 白字居中。
    color_a/b 通常用 deck.anchor / deck.comparator。"""
    dk.box(slide, cx, cy, d, d, fill=None, round=True, r=d/2,
           grad=[(0.0, color_a, 1.0), (1.0, color_b, 1.0)], grad_radial=True)
    dk.text(slide, cx, cy, d, d, [[(str(num), size, tcolor, True, False, font or dk.FONT)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

def chap(prs, deck, layout_role, num, title, sub="", accent=None):
    """章节页:编号+标题(+副标题),带半透明深色衬底条(修深暖背景白字对比坑)。
    layout_role: profile 里的章节页 role(如 'chapter');若该模板章节页背景图深暖,
    白字直接放对比不足→本 helper 自动加深色衬底。换模板若章节页是浅底,可改为深字。"""
    s = prs.slides.add_slide(prs.slide_layouts[deck.P.layout(layout_role)])
    s.placeholders[0].text = ""   # 清占位符,改自绘以控层级
    accent = accent or deck.anchor
    bx_y = 2.72
    # 半透明深色衬底条(深色取自 accent 的暗化近似)
    a = accent
    dark = RGBColor(max(0,a[0]//5), max(0,a[1]//8), max(0,a[2]//10))
    dk.box(s, 0.8, bx_y, 11.7, 1.3,
           grad=[(0.0, dark, 0.80), (1.0, dark, 0.58)], grad_angle=0, round=True, r=0.12)
    dk.text(s, 1.2, bx_y+0.18, 11.1, 0.95,
            [[(f"{num}   ", 38, GOLD_FALLBACK, True, False, dk.FONT),
              (title, 36, WHITE, True, False, dk.EAFONT)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    if sub:
        dk.text(s, 1.2, bx_y+1.42, 11.1, 0.4,
                [[(sub, 15, RGBColor(0xFF,0xE0,0xCC), False, False, dk.EAFONT)]], line_spacing=1.1)
    return s

def card(slide, x, y, w, h, fill=None, line=None, line_w=1.5, r=0.12, corners='all'):
    """圆角卡片(dk.box 的语义别名,'card' 比 'box' 好读)。"""
    return dk.box(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, round=True, corners=corners, r=r)

def para(txt, size=14, color=None, bold=False, font=None, ea=None):
    """run 元组工厂:(txt, size, color, bold, italic, font)。省记忆。"""
    return [(txt, size, color, bold, False, font or (ea or dk.EAFONT))]

def notes(slide, txt):
    """写讲者备注(讲稿)。不在幻灯片上渲染,只在演示视图显示。"""
    dk.speaker_notes(slide, txt)
