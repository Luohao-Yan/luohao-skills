#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_skeleton.py — tech-gtm-training-deck 的 build 脚手架模板(从 profile.yaml 读品牌)。

这是 examples/deepseek-harness/build.py 的通用化版本:品牌色/字体/layout 从
profile.yaml 来(非硬编),含一套完整的 deck 节奏(封面/目录/章节页/内容页/
深色关键发现/红色结论/附录)+ 动画 + 讲者备注 + lint 门禁。

用法(复制到你的工作目录改内容):
    cp templates/build_skeleton.py my-topic/build.py
    # 改 TPL=你的.pptx、逐页内容、profile.yaml 的 semantic_contract
    python my-topic/build.py

需先跑 inspect_and_profile.py 生成 profile.yaml,并填好 semantic_contract。
依赖 slide-maker skill 的 deckkit/anim。
"""
import sys, os
HERE = os.path.dirname(os.path.abspath(__file__))
# tech-gtm-training-deck/scripts(本 skill 的 scripts):本文件在 templates/,往上一级是 skill 根
SKILL_ROOT = os.path.dirname(HERE)
sys.path.insert(0, os.path.join(SKILL_ROOT, "scripts"))
from slide_maker_path import find_slide_maker
sys.path.insert(0, find_slide_maker())   # slide-maker 的 scripts(deckkit/anim,自动探测)
import deckkit as dk
from anim import Build
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from load_profile import load
from deck_helpers import (Deck, set_title, num_circle, chap, card, para, notes,
                          arch_layers, network_topo, cover, strip_branding,
                          quad_grid, steps3, code_card, text_right_card)

# === 你的模板与输出 ===
TPL = r"<改为你自己的.pptx模板路径>"
OUT = os.path.join(HERE, "training-deck.pptx")

P = load(os.path.join(HERE, "profile.yaml"))
D = Deck(P)                       # 设好 dk.FONT/EAFONT + 语义色快捷
W_IN, H_IN = D.W, D.H
WHITE = RGBColor(0xFF,0xFF,0xFF)
CARDBG = RGBColor(0xF7,0xF7,0xFA)

def build():
    prs = dk.open_template(TPL)
    strip_branding(prs)   # 清模板自带的企业 logo 图 + 版权页脚文字(build 前必做)

    # -------- 1. 封面 (cover helper:有设计的封面,非裸填占位符) --------
    cover(prs, D,
          subject="<断言式主标题:讲清是什么>",
          subtitle="<故事线:讲完听众该记住/拍板什么>",
          meta="<受众> · <日期>",
          style="band")   # band=左色带+左对齐(稳) / hero=大渐变色块+居中(愿景)
    # cover 的讲者备注(cover 不返回 slide,补 notes 用第一页)
    notes(prs.slides._sldIdLst[0].slide, "<开场:一句话讲清今天要讲什么、为什么对听众重要。>")

    # -------- 2. 目录 (content layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "目录  Catalog", D.anchor)
    col = dk.columns(2, slide=s, top=1.4, bottom=0.7, margin=0.62, gap=0.5)
    for label, items, accent in [("01 <第一部分>", ["<子项1>","<子项2>","<子项3>"], D.anchor),
                                   ("02 <第二部分>", ["<子项1>","<子项2>","<子项3>"], D.comparator)]:
        x,y,w,h = col[0] if label.startswith("01") else col[1]
        card(s, x, y, w, h, fill=CARDBG, line=accent, line_w=1.2)
        dk.text(s, x+0.3, y+0.22, w-0.6, 0.4, [[(label[:2], 22, accent, True, False, dk.FONT)]], wrap=False)
        dk.text(s, x+0.3, y+0.62, w-0.6, 0.4, [[(label[3:], 18, RGBColor(0x2A,0x2A,0x33), True, False, dk.EAFONT)]], wrap=False)
        yy = y+1.15
        for it in items:
            dk.text(s, x+0.4, yy, w-0.8, 0.32, [[("#  "+it, 13.5, RGBColor(0x2A,0x2A,0x33), False, False, dk.EAFONT)]])
            yy += 0.38
    notes(s, "<分几部分讲,为什么这个顺序。>")

    # -------- 3. 章节页 (chapter layout, 自动加深色衬底修对比) --------
    chap(prs, D, "chapter", "01", "<第一部分标题>", sub="<关键词列表>")
    notes(s, "<进入第一部分,为什么先讲这个。>")

    # -------- 4. 内容页范例 (content layout, 4-card) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<断言式标题:结论即标题>", D.anchor)
    cols = dk.columns(4, slide=s, top=1.35, bottom=1.15, margin=0.5, gap=0.28)
    data = [("<名1>","<归属1>","<形态1>","<要点1>", D.anchor),
            ("<名2>","<归属2>","<形态2>","<要点2>", D.comparator)]
    for (name,owner,form,feat,col),c in zip(data*2, cols):
        x,y,w,h = c
        card(s, x, y, w, h, fill=WHITE, line=col, line_w=1.6, r=0.1)
        dk.box(s, x, y, w, 0.52, fill=col, round=True, corners='top', r=0.1)
        dk.text(s, x+0.16, y+0.05, w-0.32, 0.42, [[(name, 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.18, y+0.64, w-0.36, 0.32, [[(owner, 13, col, True, False, dk.EAFONT)]], wrap=False)
        dk.text(s, x+0.18, y+1.0, w-0.36, 0.32, [[(form, 12.5, D.neutral, False, False, dk.EAFONT)]])
        dk.box(s, x+0.18, y+1.38, w-0.36, 0.012, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+0.18, y+1.5, w-0.36, 0.7, [[(feat, 13, RGBColor(0x2A,0x2A,0x33), False, False, dk.EAFONT)]], line_spacing=1.2)
    dk.bottom_callout(s, 0.5, W_IN-1.0, "<关键提示>", "<一句话纠偏或强调>", label_c=D.anchor, body_c=RGBColor(0x20,0x26,0x30))
    notes(s, "<这张表的讲解:逐卡点出归属/形态/要点,最后纠偏。>")

    # -------- 4b. 四宫格页型 (quad_grid:2×2 对比/分类) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<分类标题:四类是什么>", D.anchor)
    quad_grid(s, D, [
        {"tab":"<类1标签>","head":"<类1名>","body":"<类1定义与要点>"},
        {"tab":"<类2标签>","head":"<类2名>","body":"<类2定义与要点>"},
        {"tab":"<类3标签>","head":"<类3名>","body":"<类3定义与要点>"},
        {"tab":"<类4标签>","head":"<类4名>","body":"<类4定义与要点>"},
    ])
    notes(s, "<四类对比:逐格讲,红橙交替区分>")

    # -------- 4c. 三步走页型 (steps3:落地路径/阶段) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<路径标题:三步走/三阶段>", D.anchor)
    steps3(s, D, [
        {"tag":"Step 1 <阶段名>(<周期>)","head":"<阶段目标>","body":"<阶段动作要点>"},
        {"tag":"Step 2 <阶段名>(<周期>)","head":"<阶段目标>","body":"<阶段动作要点>"},
        {"tag":"Step 3 <阶段名>(<周期>)","head":"<阶段目标>","body":"<阶段动作要点>"},
    ])
    notes(s, "<三步走:按阶段顺序讲,每阶段目标+动作>")

    # -------- 4d. 左文右代码页型 (code_card:原理+代码示例) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<实现标题:某机制怎么实现>", D.anchor)
    code_card(s, D,
              left_title="<机制/原理>",
              left_body="<左栏讲机制原理,多行>",
              code_title="<快速上手代码示例>",
              code_lines=["# 注释", "import xxx", "xxx.init()"])
    notes(s, "<左讲原理右给代码,代码卡深色底>")

    # -------- 4e. 左文右大卡页型 (text_right_card:左文右架构/方案) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "<方案标题:核心场景+实现>", D.anchor)
    text_right_card(s, D,
                    left_title="<核心场景>",
                    left_body="<左栏讲场景/诉求>",
                    right_title="<实现方案/架构>",
                    right_body="<右卡放架构描述或图说明>")
    notes(s, "<左场景右方案,右卡可换 arch_layers 画架构图>")

    # -------- 5. 深色关键发现页 (dark layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("dark")])
    set_title(s, "<关键发现:断言>", WHITE)
    cols = dk.columns(2, slide=s, top=1.35, bottom=0.75, margin=0.5, gap=0.5)
    lx = cols[0]; x,y,w,h = lx
    card(s, x, y, w, h, fill=RGBColor(0x2A,0x14,0x18), line=D.anchor, line_w=1.5)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.5, [[("本机实测", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.3, y+0.6, w-0.6, 1.2, [[("<发现内容>,",16,WHITE,True,False,dk.EAFONT)]], line_spacing=1.3)
    rx = cols[1]; x,y,w,h = rx
    card(s, x, y, w, h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=1.2)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.4, [[("意味着什么", 13, D.emphasis, True, False, dk.EAFONT)]], wrap=False)
    notes(s, "<这是最该让听众记住的发现,深色页做强调。>")

    # -------- 6. 结论 (red_conclusion layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("red_conclusion")])
    set_title(s, "结论", WHITE)
    dk.text(s, 1.0, 2.2, 11.3, 2.0, [
        [("<主题> 既是", 24, WHITE, True, False, dk.EAFONT),("技术借鉴对象", 26, D.emphasis, True, False, dk.EAFONT),
         (",也是", 24, WHITE, True, False, dk.EAFONT),("战略变量", 26, D.emphasis, True, False, dk.EAFONT),("。", 24, WHITE, True, False, dk.EAFONT)],
        [("", 8, WHITE, False, False)],
        [("短期:<...>", 16, WHITE, False, False, dk.EAFONT)],
        [("", 6, WHITE, False, False)],
        [("中期:<...>", 16, WHITE, False, False, dk.EAFONT)],
        [("", 6, WHITE, False, False)],
        [("长期:<...>", 16, WHITE, False, False, dk.EAFONT)],
    ], line_spacing=1.2)
    notes(s, "<收尾:一句话点睛 + 短中长期行动。>")

    # -------- 7. 附录 (content layout) --------
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    set_title(s, "附录·证据出处", D.anchor)
    notes(s, "<证据出处:按来源分组列 file_path:line。>")

    # -------- (可选) 分层架构图页 / 网络拓扑图页 --------
    # brief.need_arch_diagram / brief.tilt=tech 时画一页架构图:
    #   s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    #   set_title(s, "技术架构", D.anchor)
    #   arch_layers(s, [
    #       {"name":"接入层","items":["Web","App"],"height":0.8},
    #       {"name":"服务层","items":["svcA","svcB","svcC"],"height":1.4},
    #       {"name":"存储层","items":["DB"],"height":0.7},
    #   ], x=0.4, y=1.4, w=12.5, total_h=4.6, accent=D.neutral)
    #   notes(s, "分层架构:从上到下 接入→服务→存储。")
    # brief.need_network_topo 或 doc 含网络/部署拓扑 时画一页拓扑:
    #   s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")])
    #   set_title(s, "网络拓扑", D.anchor)
    #   network_topo(s,
    #     nodes=[{"id":"net","kind":"cloud","x":0.5,"y":0.1,"label":"Internet"},
    #            {"id":"fw","kind":"firewall","x":0.5,"y":0.4,"label":"防火墙"},
    #            {"id":"sw","kind":"switch","x":0.5,"y":0.7,"label":"交换机"},
    #            {"id":"s1","kind":"server","x":0.2,"y":0.95,"label":"应用"}],
    #     links=[{"from":"net","to":"fw","label":"专线"},{"from":"fw","to":"sw"},
    #            {"from":"sw","to":"s1","label":"千兆"}],
    #     accent=D.neutral)
    #   notes(s, "网络拓扑:Internet→防火墙→交换机→服务器。")

    dk.lint_layout(prs, strict=True)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
