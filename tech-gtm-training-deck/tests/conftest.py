# -*- coding: utf-8 -*-
"""共享测试夹具:注入 slide-maker 路径 + 造一个空 deckkit-friendly Presentation。"""
import os, sys
HERE = os.path.dirname(os.path.abspath(__file__))
SKILL_SCRIPTS = os.path.join(os.path.dirname(HERE), "scripts")
sys.path.insert(0, SKILL_SCRIPTS)
from slide_maker_path import find_slide_maker
sys.path.insert(0, find_slide_maker())
import deckkit as dk  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# 一个极简模板:13.333x7.5in, 一个 content layout(只标题占位符)
def make_test_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    """加一张空白页(无 layout 占位符干扰),helper 直接在自绘几何上画。"""
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank
