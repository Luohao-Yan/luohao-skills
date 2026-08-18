#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""inspect_and_profile.py — 探查用户 .pptx 模板,生成 profile.yaml + profile.md。

这是 Stage 3 的第一步:把用户的企业模板沉淀成机器可读的品牌档案,让 build 脚本
从 profile.yaml 读配色/字体/layout idx,而非人工抄硬编码(漂移坑)。

用法:
    python inspect_and_profile.py <user.pptx> [--out <dir>]
    # 生成 <out>/profile.yaml (机器读) + <out>/profile.md (人读草稿)
    # 默认 --out 为当前目录

它用 python-pptx 读:slide size、layouts(idx+name+占位符几何)、theme clrScheme、
font scheme。不依赖 slide-maker 的 inspect_template(独立可跑);但若 slide-maker
已装,可对比印证。语义色契约(semantic_contract)需人工填——accent→本deck语义角色
与内容绑定,机器无法推断。
"""
import sys, os, re, zipfile, argparse
from pptx import Presentation
from pptx.util import Emu

def emu_in(v):
    return round(v / 914400, 3) if v is not None else None

def extract_theme(pptx_path):
    """从 ppt/theme/theme1.xml 提取 clrScheme + fontScheme。"""
    z = zipfile.ZipFile(pptx_path)
    try:
        theme = z.read("ppt/theme/theme1.xml").decode("utf-8")
    except KeyError:
        return {}, {}
    colors = {}
    m = re.search(r"<a:clrScheme.*?</a:clrScheme>", theme, re.S)
    if m:
        scheme = m.group(0)
        for name in ["dk1","lt1","dk2","lt2","accent1","accent2","accent3","accent4","accent5","accent6","hlink"]:
            mm = re.search(r"<a:%s>(.*?)</a:%s>" % (name, name), scheme, re.S)
            if mm:
                c = re.search(r'(?:srgbClr|sysClr) (?:val|lastClr)="([0-9A-Fa-f]{6})"', mm.group(1))
                if c:
                    colors[name] = "#" + c.group(1).upper()
    fonts = {}
    m2 = re.search(r"<a:fontScheme.*?</a:fontScheme>", theme, re.S)
    if m2:
        fs = m2.group(0)
        for role in ["majorFont","minorFont"]:
            mm = re.search(r"<a:%s>(.*?)</a:%s>" % (role, role), fs, re.S)
            if mm:
                latin = re.search(r'<a:latin typeface="([^"]*)"', mm.group(1))
                ea = re.search(r'<a:ea typeface="([^"]*)"', mm.group(1))
                fonts[role] = {"latin": latin.group(1) if latin else "", "ea": ea.group(1) if ea else ""}
    return colors, fonts

def inspect_layouts(prs):
    layouts = {}
    for i, lay in enumerate(prs.slide_layouts):
        phs = []
        for ph in lay.placeholders:
            phs.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "x": emu_in(ph.left), "y": emu_in(ph.top),
                "w": emu_in(ph.width), "h": emu_in(ph.height),
            })
        layouts[i] = {"name": lay.name, "placeholders": phs}
    return layouts

def pick_key_layouts(layouts):
    """从 layouts 里挑出 build 脚本常用的角色:cover/content/dark/red/chapter/blank。
    返回 role -> idx 映射(给 profile.yaml 的 layouts.key 段用)。启发式按 name 关键字。"""
    role_map = {}
    for idx, info in layouts.items():
        n = info["name"]
        if any(k in n for k in ["标题幻灯片","标题页","Title Slide"]) and "cover" not in role_map:
            role_map["cover"] = idx
        elif "章节" in n and "chapter" not in role_map:
            role_map["chapter"] = idx
        elif "深色" in n and "有底线" in n and "dark" not in role_map:
            role_map["dark"] = idx
        elif "红色" in n and "有底线" in n and "red_conclusion" not in role_map:
            role_map["red_conclusion"] = idx
        elif "浅色" in n and "有底线" in n and "content" not in role_map:
            role_map["content"] = idx
        elif n.strip() == "空白" and "blank" not in role_map:
            role_map["blank"] = idx
    return role_map

def write_yaml(pptx_path, colors, fonts, layouts, key_layouts, w_in, h_in, out_yaml):
    lines = [
        "# profile.yaml — 由 inspect_and_profile.py 生成,build 脚本通过 load_profile.py 加载。",
        "# 不要手改 colors/fonts/layouts(会漂移);只填 semantic_contract(语义角色绑定)。",
        f"# source: {pptx_path}",
        "",
        "canvas:",
        f"  w_in: {w_in}",
        f"  h_in: {h_in}",
        "",
        "colors:",
    ]
    for k, v in colors.items():
        lines.append(f"  {k}: \"{v}\"")
    lines += ["", "fonts:"]
    mf = fonts.get("majorFont", {})
    lines.append(f"  latin: \"{mf.get('latin','')}\"   # majorFont latin")
    lines.append(f"  ea: \"{mf.get('ea','')}\"          # majorFont ea (CJK); build 用它打 <a:ea> tag")
    lines += ["", "# 常用渐进三色(人工从 accent 选,供流程图渐进色用;可选):",
              "  gradient3: []   # e.g. [\"#FCE5D6\", \"#F46B30\", \"#C60000\"]"]
    lines += ["", "layouts:"]
    for role, idx in key_layouts.items():
        lines.append(f"  {role}: {idx}   # {layouts[idx]['name']}")
    lines += ["", "# 语义色契约(人工填):把 accent 绑到本 deck 的语义角色。与内容绑定,机器无法推断。",
              "# 例: anchor_subject -> accent1, comparator -> accent2, neutral -> accent3, emphasis -> accent4",
              "semantic_contract: {}"]
    open(out_yaml, "w", encoding="utf-8").write("\n".join(lines) + "\n")
    return out_yaml

def write_md(pptx_path, colors, fonts, layouts, key_layouts, w_in, h_in, out_md):
    lines = [
        f"# 模板品牌 profile · {os.path.basename(pptx_path)}",
        "",
        f"> 由 `inspect_and_profile.py` 生成草稿。`colors`/`fonts`/`layouts` 自动;**语义色契约、视觉语汇、footprint 需人工补**。",
        "",
        "## Canvas",
        f"- {round(w_in/914400*1,0) if False else ''}{w_in} × {h_in} in",
        "",
        "## 主题色 (theme clrScheme)",
    ]
    for k, v in colors.items():
        lines.append(f"- **{k} {v}** — (填语义角色,如:主色/核心结论)")
    lines += ["", "## 字体",
              f"- majorFont latin=`{fonts.get('majorFont',{}).get('latin','')}` ea=`{fonts.get('majorFont',{}).get('ea','')}`",
              f"- minorFont latin=`{fonts.get('minorFont',{}).get('latin','')}` ea=`{fonts.get('minorFont',{}).get('ea','')}`",
              "- deckkit 设置: `FONT=latin`, `EAFONT=ea`(CJK), `DISPLAY=latin`",
              "", "## 版式 (layouts, build 脚本按 role 取 idx)"]
    for role, idx in key_layouts.items():
        lines.append(f"- layout {idx} `{layouts[idx]['name']}` → role: **{role}**")
    lines += ["", "## 语义色契约 (本 deck,人工填)",
              "- (把上面 accent 绑到本 deck 语义角色:anchor_subject / comparator / neutral / emphasis)",
              "", "## 内容页视觉语汇 (人工补,沿用模板)",
              "- 卡片 / 编号圆 / 箭头 / 底线 / 强调句 等模板自带语汇",
              "", "## 注意 / footprint",
              "- 内容页占位符仅标题;正文用 deckkit 画,避开版式自带底线(底部约 0.5in)",
              "- 章节页若背景图深暖,白字对比不足 → 加半透明深色衬底条(deck_helpers.chap)",
              "- 本机若无 SVG 光栅器 → 不用 SVG 图标,沿用模板几何语汇"]
    open(out_md, "w", encoding="utf-8").write("\n".join(lines) + "\n")
    return out_md

def main():
    ap = argparse.ArgumentParser(description="探查 .pptx 模板 → 生成 profile.yaml + profile.md")
    ap.add_argument("pptx", help="用户 .pptx 模板路径")
    ap.add_argument("--out", default=".", help="输出目录(默认当前)")
    args = ap.parse_args()
    if not os.path.isfile(args.pptx):
        sys.exit(f"not found: {args.pptx}")
    prs = Presentation(args.pptx)
    w_in, h_in = emu_in(prs.slide_width), emu_in(prs.slide_height)
    colors, fonts = extract_theme(args.pptx)
    layouts = inspect_layouts(prs)
    key_layouts = pick_key_layouts(layouts)
    os.makedirs(args.out, exist_ok=True)
    y = write_yaml(args.pptx, colors, fonts, layouts, key_layouts, w_in, h_in, os.path.join(args.out, "profile.yaml"))
    m = write_md(args.pptx, colors, fonts, layouts, key_layouts, w_in, h_in, os.path.join(args.out, "profile.md"))
    print(f"canvas: {w_in} x {h_in} in")
    print(f"colors: {len(colors)} accents, fonts: latin={fonts.get('majorFont',{}).get('latin','')} ea={fonts.get('majorFont',{}).get('ea','')}")
    print(f"layouts: {len(layouts)} total, key roles: {key_layouts}")
    print(f"-> {y}")
    print(f"-> {m}")
    print("next: 填 profile.yaml 的 semantic_contract,然后用 load_profile.py 在 build 脚本里加载")

if __name__ == "__main__":
    main()
