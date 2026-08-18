#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""deck_helpers.py — 在模板上构建 deck 的通用积木,颜色/字体从 profile 来(非硬编)。

build 脚本:
    import sys, os
    sys.path.insert(0, find_slide_maker())          # deckkit/anim 在这(自动探测)
    sys.path.insert(0, "<tech-gtm-training-deck>/scripts")
    import deckkit as dk
    from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes

这些 helper 把"在用户模板上画一页"的重复动作封装好:填标题占位符并打 CJK 字体
tag、画模板语汇的渐变编号圆、画带半透明衬底的章节页(修深暖背景白字对比坑)、
画卡片、写讲者备注。颜色都从传入的 Profile 来,不硬编 hex。
"""
import os, sys
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import deckkit as dk

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD_FALLBACK = RGBColor(0xFF, 0xC0, 0x00)
INK = RGBColor(0x2A, 0x2A, 0x33)
MUTE = RGBColor(0x6B, 0x6B, 0x6B)

_IN = 914400  # EMU per inch, for shape position judgement

def _in(v):
    """Emu -> inch (None/0 safe)."""
    return (v or 0) / _IN

def bottom_callout_at(slide, x, w, bottom_y, label, body, **kw):
    """锚定到指定 bottom_y(向上长)的 callout,精确控制贴底位置——比 bottom_callout 的
    footer 锚更可控(后者受 FOOTER_BAND=0.5 限制,底到 6.85 留白偏大)。
    先量高度,再放 y=bottom_y-h。返回 top y(供上方内容避让)。
    用法: top = bottom_callout_at(s, 0.5, W-1.0, 7.07, "一句话", "...");
          上方内容 bottom 应 < top - 0.28(留 0.28in 视觉间距)。"""
    ch = dk.measure_callout(label, body, w)
    y = bottom_y - ch
    dk.callout(slide, x, y, w, ch, label, body, **kw)
    return y

class Deck:
    """持有 profile + deckkit 全局字体,build 脚本从这取颜色/字体/layout。"""
    def __init__(self, profile):
        self.P = profile
        latin, ea = profile.fonts()
        dk.FONT = latin; dk.EAFONT = ea; dk.DISPLAY = latin
        self.W, self.H = profile.canvas()
    # 语义色快捷(走 profile 的 semantic_contract)
    @property
    def anchor(self): return self.P.color("anchor_subject")
    @property
    def comparator(self): return self.P.color("comparator")
    @property
    def neutral(self): return self.P.color("neutral")
    @property
    def emphasis(self): return self.P.color("emphasis")

def set_title(slide, text, color, size=26, bold=True, font=None, ea=None):
    """填 idx=0 标题占位符 + 设字号/粗/色 + 打 CJK ea 字体 tag(kinsoku/渲染才正确)。
    每个内容页都要做;颜色由调用方从 profile 传(通常是 anchor 主色)。"""
    ph = slide.placeholders[0]
    ph.text = text
    latin, eaf = (font or dk.FONT), (ea or dk.EAFONT)
    for p in ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = latin
            dk._apply_ea(r, eaf)
    return ph

def num_circle(slide, cx, cy, d, num, color_a, color_b, tcolor=WHITE, size=16, font=None):
    """渐变编号圆(模板语汇):color_a→color_b 径向渐变 + 白字居中。
    color_a/b 通常用 deck.anchor / deck.comparator。"""
    dk.box(slide, cx, cy, d, d, fill=None, round=True, r=d/2,
           grad=[(0.0, color_a, 1.0), (1.0, color_b, 1.0)], grad_radial=True)
    dk.text(slide, cx, cy, d, d, [[(str(num), size, tcolor, True, False, font or dk.FONT)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

def chap(prs, deck, layout_role, num, title, sub="", accent=None):
    """章节页:用模板的章节页占位符(标题 idx0 + 副标题 idx10),字号/位置/背景由版式保证。
    不自绘衬底——沿用模板设计语言(模板的占位符已在合适位置,背景图标题区是暗的,白字可读)。
    若某模板章节页白字确不可读,再说;默认信任模板设计。"""
    accent = accent or deck.anchor
    s = prs.slides.add_slide(prs.slide_layouts[deck.P.layout(layout_role)])
    # idx0 标题占位符:编号 + 标题
    ph = s.placeholders[0]
    ph.text = f"{num}  {title}"
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.name = dk.FONT
            dk._apply_ea(r, dk.EAFONT)
    # idx10 副标题占位符(模板自带 24 号样式)
    if sub:
        try:
            ph2 = s.placeholders[10]; ph2.text = sub
            for p in ph2.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = dk.EAFONT; dk._apply_ea(r, dk.EAFONT)
        except (KeyError, IndexError):
            pass
    return s

def card(slide, x, y, w, h, fill=None, line=None, line_w=1.5, r=0.12, corners='all'):
    """圆角卡片(dk.box 的语义别名,'card' 比 'box' 好读)。"""
    return dk.box(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, round=True, corners=corners, r=r)

def para(txt, size=14, color=None, bold=False, font=None, ea=None):
    """run 元组工厂:(txt, size, color, bold, italic, font)。省记忆。"""
    return [(txt, size, color, bold, False, font or (ea or dk.EAFONT))]

def notes(slide, txt):
    """写讲者备注(讲稿)。不在幻灯片上渲染,只在演示视图显示。"""
    dk.speaker_notes(slide, txt)


# ---- 分层架构图 (arch_layers) ----
# 复刻珠海项目孵化 PPT p05 画法:全宽色带分层 + 层内组件块 + 交替浅色
ARCH_TINTS = [  # colored 风格的 5 色浅色梯度循环
    RGBColor(0xEF, 0xF6, 0xFF),  # 浅蓝
    RGBColor(0xF0, 0xFD, 0xF4),  # 浅绿
    RGBColor(0xFF, 0xF7, 0xED),  # 浅橙
    RGBColor(0xED, 0xE9, 0xFE),  # 浅紫
    RGBColor(0xCC, 0xFB, 0xF1),  # 浅青
]
ARCH_MONO_TINT = RGBColor(0xF3, 0xF4, 0xF6)
ARCH_COMPONENT_FILL = RGBColor(0xFF, 0xFF, 0xFF)  # 组件块白底
ARCH_DARK_CAP = RGBColor(0x1F, 0x29, 0x37)        # 末层深色收尾(可选)

def arch_layers(slide, layers, x=0.4, y=1.4, w=12.5, total_h=5.2,
                style="colored", accent=None, font=None):
    """分层架构图。layers=[{name,items:[...],height,color?},...]。
    style: colored(多彩分层)/ mono(单色分层)。
    返回各层中心点 [(cx,cy),...](供后续连线)。
    画法:每层一条全宽色带(L=x,W=w,浅色),层名左上,层内组件块横向并排(白底+accent描边)。
    超过6层:自动压缩高度并 print 提示建议拆页。"""
    n = len(layers)
    if n == 0:
        return []
    if n > 6:
        print("[arch_layers] 层数较多(>6),建议拆页;已自动压缩层高")
    sum_h = sum(L.get("height", 0.8) for L in layers)
    if sum_h > total_h:
        scale = total_h / sum_h
    else:
        scale = 1.0
    accent = accent or RGBColor(0x3F, 0x54, 0x69)
    ink = RGBColor(0x2A, 0x2A, 0x33)
    ea = font or dk.EAFONT
    cy_list = []
    cur_y = y
    for i, L in enumerate(layers):
        h = L.get("height", 0.8) * scale
        tint = L.get("color") or (ARCH_TINTS[i % len(ARCH_TINTS)] if style == "colored" else ARCH_MONO_TINT)
        # 全宽色带
        dk.box(slide, x, cur_y, w, h, fill=tint, round=True, r=0.06)
        # 层名(左上)
        dk.text(slide, x + 0.16, cur_y + 0.04, 2.4, 0.3,
                [[(L.get("name", ""), 12, RGBColor(0x4B,0x55,0x63), True, False, ea)]], wrap=False)
        # 层内组件块:横向均分
        items = L.get("items", [])
        if items:
            pad = 0.16
            inner_x = x + 2.6
            inner_w = w - 2.6 - pad
            gap = 0.12
            cw = (inner_w - gap * (len(items) - 1)) / len(items)
            ch = min(0.28, h - 0.16)
            comp_y = cur_y + (h - ch) / 2
            for j, it in enumerate(items):
                cx = inner_x + j * (cw + gap)
                dk.box(slide, cx, comp_y, cw, ch, fill=ARCH_COMPONENT_FILL,
                       line=accent, line_w=1.0, round=True, r=0.05)
                dk.text(slide, cx, comp_y, cw, ch,
                        [[(it, 10.5, ink, False, False, ea)]],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        cy_list.append((x + w / 2, cur_y + h / 2))
        cur_y += h
    return cy_list


# ---- 网络拓扑图 (network_topo) ----
def _icon_path(kind):
    """assets/icons/<kind>.png 的绝对路径(相对本 skill 根)。不存在返回 None。"""
    here = os.path.dirname(os.path.abspath(__file__))
    p = os.path.join(os.path.dirname(here), "assets", "icons", f"{kind}.png")
    return p if os.path.isfile(p) else None

def network_topo(slide, nodes, links, x=0.5, y=1.4, w=12.3, h=5.3,
                 accent=None, font=None, icon_size=0.5):
    """网络拓扑图。nodes=[{id,kind,x,y,label,sub?}] (x,y 为 0..1 相对坐标);
    links=[{from,to,label?,style?,arrow?}]。
    用 deckkit connect_boxes 边到边精准连线(不穿节点);图标从 assets/icons/<kind>.png 读,
    缺失则降级为形状节点(圆角矩形+kind 文字)。颜色/字体从 accent/font 来。
    返回 {id: ((cx,cy), rect)} 供外部追加连线。"""
    accent = accent or RGBColor(0x3F, 0x54, 0x69)
    ea = font or dk.EAFONT
    ink = RGBColor(0x2A, 0x2A, 0x33)
    # 节点尺寸(图标 + 标签)
    nw, nh = icon_size + 0.6, icon_size + 0.5
    rects = {}
    # 1. 先算所有节点 rect(供连线),再画连线(画在节点下),再画节点(盖上)
    for nd in nodes:
        nx = x + nd["x"] * (w - nw)
        ny = y + nd["y"] * (h - nh)
        rects[nd["id"]] = (nx, ny, nw, nh)
    # 2. 连线(z-order: 先连线)
    for lk in links:
        a = rects[lk["from"]]; b = rects[lk["to"]]
        kw = {"style": lk.get("style","solid"), "color": accent, "width": 1.4,
              "label": lk.get("label",""), "arrow": lk.get("arrow", True)}
        dk.connect_boxes(slide, a, b, **kw)
    # 3. 节点(z-order: 后画,盖住连线 seam)
    for nd in nodes:
        nx, ny, nw2, nh2 = rects[nd["id"]]
        ip = _icon_path(nd["kind"])
        if ip:
            dk.icon(slide, ip, nx + (nw2 - icon_size)/2, ny + 0.06, icon_size, disc=accent)
        else:
            # 降级:圆角矩形 + kind 文字
            dk.box(slide, nx, ny, nw2, nh2, fill=RGBColor(0xFF,0xFF,0xFF),
                   line=accent, line_w=1.2, round=True, r=0.08)
            dk.text(slide, nx, ny+0.02, nw2, 0.3,
                    [[(nd["kind"], 9, accent, True, False, ea)]],
                    align=PP_ALIGN.CENTER, wrap=False)
        # 标签(图标下方)
        dk.text(slide, nx, ny + nh2 - 0.26, nw2, 0.24,
                [[(nd.get("label",""), 10, ink, True, False, ea)]],
                align=PP_ALIGN.CENTER, wrap=False)
    return {nid: ((r[0]+r[2]/2, r[1]+r[3]/2), r) for nid, r in rects.items()}


# ---- 封面 (cover) ----
def cover(prs, deck, subject, subtitle="", meta="", style="band"):
    """有设计的封面(非裸填模板占位符)。在 blank layout 上自绘,避免继承模板的 logo 版式。

    subject : 主标题(断言式主题句,非裸主题词)
    subtitle: 故事线/副标题(讲完听众该记住或拍板什么),1 行
    meta    : 受众 + 日期一行(如"技术团队 · 2026/08/19")
    style   : "band"(左色带 + 左对齐,稳,默认) / "hero"(大色块 + 居中,适合愿景)
    颜色全走 deck.anchor/comparator/neutral(profile),不硬编 hex。
    返回 slide。
    """
    s = prs.slides.add_slide(prs.slide_layouts[deck.P.layout("blank")])
    W, H = deck.W, deck.H
    anchor = deck.anchor
    neutral = deck.neutral
    ink = INK
    ea = dk.EAFONT
    latin = dk.FONT
    if style == "hero":
        # 上半渐变色块(主色深浅渐变,学模板的纵向高光)+ 居中标题
        dk.box(s, 0, 0, W, H*0.58,
               grad=[(0.0, anchor, 1.0), (1.0, deck.comparator, 0.85)], grad_angle=90, round=False)
        dk.text(s, 0.6, H*0.16, W-1.2, 1.4,
                [[(subject, 40, WHITE, True, False, ea)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        if subtitle:
            dk.text(s, 0.6, H*0.16+1.5, W-1.2, 0.5,
                    [[(subtitle, 16, RGBColor(0xF3,0xF5,0xF8), False, False, ea)]],
                    align=PP_ALIGN.CENTER)
        if meta:
            dk.text(s, 0.6, H*0.62, W-1.2, 0.4,
                    [[(meta, 13, neutral, False, False, ea)]], align=PP_ALIGN.CENTER)
        # 底部 comparator 渐变细带收尾
        dk.box(s, 0, H-0.18, W, 0.18,
               grad=[(0.0, deck.comparator, 1.0), (1.0, anchor, 1.0)], grad_angle=0, round=False)
    else:
        # band: 左侧 anchor 渐变色带(主色深浅)+ 左对齐标题
        dk.box(s, 0, 0, 0.32, H,
               grad=[(0.0, anchor, 1.0), (1.0, deck.comparator, 0.9)], grad_angle=90, round=False)
        # 主标题(垂直偏上,左对齐,顶部对齐)
        dk.text(s, 1.1, H*0.28, W-1.6, 1.3,
                [[(subject, 38, anchor, True, False, ea)]],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.12, wrap=True)
        # 副标题/故事线(紧跟主标题下方)
        if subtitle:
            dk.text(s, 1.12, H*0.28+1.35, W-1.7, 0.9,
                    [[(subtitle, 16, neutral, False, False, ea)]],
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25, wrap=True)
        # meta 行 + 分隔细线(贴底)
        if meta:
            dk.box(s, 1.12, H-0.95, 1.6, 0.02, fill=anchor)
            dk.text(s, 1.12, H-0.80, W-1.7, 0.35,
                    [[(meta, 13, MUTE, False, False, ea)]], align=PP_ALIGN.LEFT)
    return s


# ---- 清模板品牌 logo/页脚 (strip_branding) ----
# 模板各 layout 右上角常带企业 logo 图(约 1.0x0.32in,在画布右上角区);
# 部分 layout 还有"北京金山云网络技术有限公司 / WWW.KSYUN.COM / Copyright"版权页脚文字。
# open_template 保留母版/版式,这些会页页继承——必须 build 前清掉。
BRAND_TEXT_KEYS = ("金山云", "KSYUN", "Copyright", "北京金山云网络技术", "金山软件", "kingsoft", "ksyun")

def _is_logo_pic(sh, W_in, H_in):
    """判定一个 picture shape 是不是右上角 logo:位置在右上角区(left>W*0.79 且 top<H*0.13)
    且尺寸小(w<1.6in 且 h<0.7in)。大背景图(如章节页装饰)不在此区,保留。"""
    try:
        if str(sh.shape_type) != "PICTURE (13)":
            return False
        l, t = _in(sh.left), _in(sh.top)
        w, h = _in(sh.width), _in(sh.height)
        return l > W_in * 0.79 and t < H_in * 0.13 and w < 1.6 and h < 0.7
    except Exception:
        return False

def _is_brand_text(sh):
    """判定一个文本框是否含品牌版权文字(页脚)。"""
    try:
        if not sh.has_text_frame:
            return False
        t = sh.text_frame.text or ""
        low = t.lower()
        return any(k.lower() in low for k in BRAND_TEXT_KEYS)
    except Exception:
        return False

def strip_branding(prs, keep_logo=False, verbose=False):
    """在 open_template 后、add_slide 前 调一次。清掉母版+各版式上的企业 logo 图和
    版权页脚文字(含金山云等)。keep_logo=True 则跳过(留给真要保留品牌的对外汇报场景)。
    返回删除计数 {pics, texts}。章节页等大装饰背景图(非右上角小图)保留。"""
    if keep_logo:
        return {"pics": 0, "texts": 0}
    W_in = prs.slide_width / _IN
    H_in = prs.slide_height / _IN
    npic = ntxt = 0
    targets = list(prs.slide_masters) + list(prs.slide_layouts)
    for host in targets:
        # 收集要删的 shape(先收集再删,避免遍历中改集合)
        to_del = []
        for sh in list(host.shapes):
            if _is_logo_pic(sh, W_in, H_in):
                to_del.append(sh); npic += 1
                continue
            if _is_brand_text(sh):
                to_del.append(sh); ntxt += 1
        for sh in to_del:
            sh._element.getparent().remove(sh._element)
    if verbose:
        print(f"[strip_branding] removed {npic} logo pic(s), {ntxt} brand text(s)")
    return {"pics": npic, "texts": ntxt}


# ---- 内容页型 helpers (抽自模板页骨架,配色走 profile) ----

def quad_grid(slide, deck, items, x=0.8, y=1.5, w=12.1, total_h=4.9,
              gap=0.3, accent=None, font=None):
    """2×2 四宫格(学记忆模板 slide7 同款)。items=[{tab,head,body,color?},...] 共 4 项。
    每卡:顶部 tab 角色标签 + 标题 + 正文。配色用 deck.anchor/comparator 交替或 items 自带 color。"""
    ea = font or dk.EAFONT
    ink = INK
    anchor = accent or deck.anchor
    comp = deck.comparator
    cw = (w - gap) / 2
    ch = (total_h - gap) / 2
    rects = []
    palette_cycle = [anchor, comp, comp, anchor]
    for i, it in enumerate(items):
        col = it.get("color", palette_cycle[i % 4])
        cx = x + (i % 2) * (cw + gap)
        cy = y + (i // 2) * (ch + gap)
        card(slide, cx, cy, cw, ch, fill=WHITE, line=col, line_w=1.6, r=0.1)
        dk.box(slide, cx, cy, cw, 0.06, fill=col, round=True, corners='top', r=0.1)
        dk.text(slide, cx+0.3, cy+0.18, cw-0.6, 0.4,
                [[(it.get("tab",""), 13.5, col, True, False, ea)]], wrap=False)
        dk.text(slide, cx+0.3, cy+0.6, cw-0.6, 0.4,
                [[(it.get("head",""), 15, ink, True, False, ea)]], wrap=False)
        dk.box(slide, cx+0.3, cy+1.02, cw-0.6, 0.012, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(slide, cx+0.3, cy+1.12, cw-0.6, ch-1.2,
                [[(it.get("body",""), 12.5, ink, False, False, ea)]], line_spacing=1.25, wrap=True)
        rects.append((cx, cy, cw, ch))
    return rects

def steps3(slide, deck, steps, x=0.7, y=1.3, w=12.1, h=5.0, gap=0.3, accent=None, font=None):
    """三步走纵列(学记忆模板 slide22 同款)。steps=[{tag,head,body?,color?},...] 共 3 项。
    每列:顶部 Step N 标签条(色块)+ 大卡 + 正文。"""
    ea = font or dk.EAFONT
    ink = INK
    anchor = accent or deck.anchor
    comp = deck.comparator
    cw = (w - gap * 2) / 3
    rects = []
    cols = [anchor, comp, deck.neutral]
    for i, st in enumerate(steps):
        col = st.get("color", cols[i % 3])
        cx = x + i * (cw + gap)
        dk.box(slide, cx, y, cw, 0.6, fill=col, round=True, corners='top', r=0.1)
        dk.text(slide, cx+0.25, y+0.06, cw-0.5, 0.5,
                [[(st.get("tag", f"Step {i+1}"), 14, WHITE, True, False, ea)]],
                anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.box(slide, cx, y+0.6, cw, h-0.6, fill=WHITE, line=col, line_w=1.4, round=True, corners='bottom', r=0.1)
        dk.text(slide, cx+0.3, y+0.9, cw-0.6, 0.5,
                [[(st.get("head",""), 16, col, True, False, ea)]], wrap=False)
        if st.get("body"):
            dk.text(slide, cx+0.3, y+1.45, cw-0.6, h-1.6,
                    [[(st["body"], 13, ink, False, False, ea)]], line_spacing=1.3, wrap=True)
        rects.append((cx, y, cw, h))
    return rects

def code_card(slide, deck, left_title, left_body, code_title, code_lines, x=0.6, y=1.1,
              left_w=6.0, right_w=5.7, h=6.2, accent=None, font=None):
    """左正文 + 右代码大卡(学记忆模板 slide16 同款)。
    left_title: 左栏小标题; left_body: 左栏正文; code_title: 右卡标题;
    code_lines: 右卡代码(list[str] 或 str)。"""
    ea = font or dk.EAFONT
    ink = INK
    anchor = accent or deck.anchor
    dk.text(slide, x, y, left_w, 0.4,
            [[(left_title, 14, anchor, True, False, ea)]], wrap=False)
    dk.text(slide, x, y+0.5, left_w, h-0.5,
            [[(left_body, 14, ink, False, False, ea)]], line_spacing=1.4, wrap=True)
    rx = x + left_w + 0.3
    card(slide, rx, y, right_w, h, fill=RGBColor(0x1F,0x29,0x37), line=anchor, line_w=1.2, r=0.1)
    dk.text(slide, rx+0.3, y+0.15, right_w-0.6, 0.4,
            [[(code_title, 14, anchor, True, False, ea)]], wrap=False)
    dk.box(slide, rx+0.3, y+0.58, right_w-0.6, 0.02, fill=anchor)
    code_str = "\n".join(code_lines) if isinstance(code_lines, list) else code_lines
    dk.text(slide, rx+0.3, y+0.72, right_w-0.6, h-0.9,
            [[(code_str, 11.5, RGBColor(0xA9,0xB7,0xCA), False, False, ea)]],
            line_spacing=1.25, wrap=True)

def text_right_card(slide, deck, left_title, left_body, right_title, right_body, x=0.6, y=1.1,
                    left_w=6.0, right_w=5.7, h=6.2, accent=None, font=None):
    """左正文 + 右大卡(学记忆模板 slide17/21 同款)。右卡放架构/方案描述。"""
    ea = font or dk.EAFONT
    ink = INK
    anchor = accent or deck.anchor
    dk.text(slide, x, y, left_w, 0.4,
            [[(left_title, 14, anchor, True, False, ea)]], wrap=False)
    dk.text(slide, x, y+0.5, left_w, h-0.5,
            [[(left_body, 14, ink, False, False, ea)]], line_spacing=1.4, wrap=True)
    rx = x + left_w + 0.3
    card(slide, rx, y, right_w, h, fill=WHITE, line=anchor, line_w=1.4, r=0.1)
    dk.text(slide, rx+0.3, y+0.18, right_w-0.6, 0.4,
            [[(right_title, 14, anchor, True, False, ea)]], wrap=False)
    dk.box(slide, rx+0.3, y+0.6, right_w-0.6, 0.02, fill=anchor)
    dk.text(slide, rx+0.3, y+0.75, right_w-0.6, h-0.95,
            [[(right_body, 13, ink, False, False, ea)]], line_spacing=1.35, wrap=True)
