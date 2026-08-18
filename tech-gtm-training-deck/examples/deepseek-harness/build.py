# -*- coding: utf-8 -*-
"""build.py — DeepSeek Harness 培训 deck(示例,品牌色从 profile.yaml 读,非硬编)

tech-gtm-training-deck skill 的 Stage 3 产物示例。复用 slide-maker 的 deckkit/anim,
品牌色/字体/layout 从 profile.yaml 通过 deck_helpers.Deck 加载。

跑通需:① 装好 slide-maker skill;② 把 TPL 改成你的 .pptx;③ profile.yaml 填好 semantic_contract。
"""
import sys, os
HERE = os.path.dirname(os.path.abspath(__file__))   # examples/deepseek-harness/
SKILL_ROOT = os.path.dirname(os.path.dirname(HERE))  # tech-gtm-training-deck/
sys.path.insert(0, os.path.join(SKILL_ROOT, "scripts"))   # 本 skill 的 scripts(含 slide_maker_path/load_profile/deck_helpers)
from slide_maker_path import find_slide_maker
sys.path.insert(0, find_slide_maker())   # slide-maker 的 scripts(deckkit/anim,自动探测)
import deckkit as dk
from anim import Build
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from load_profile import load
from deck_helpers import Deck, set_title, num_circle, chap, card, para, notes, bottom_callout_at

TPL = r"<改为你自己的.pptx模板路径>"   # 例:r"C:/Users/.../某企业模板.pptx"
OUT = os.path.join(HERE, "training-deck.pptx")

P = load(os.path.join(HERE, "profile.yaml"))
D = Deck(P)
W_IN, H_IN = D.W, D.H
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2A,0x2A,0x33)
INK = RGBColor(0x2A,0x2A,0x33)
MUTE = RGBColor(0x6B,0x6B,0x6B)
CARDBG = RGBColor(0xF7,0xF7,0xFA)
DEEP = RGBColor(0xC6,0x00,0x00)

def build():
    prs = dk.open_template(TPL)

    # 1 封面
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("cover")])
    try:
        t = s.placeholders[0]; t.text = "DeepSeek Harness 能力培训"
        for p in t.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(40); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=dk.EAFONT
                dk._apply_ea(r, dk.EAFONT)
    except Exception: pass
    for pidx, txt in [(10,"三问三追问 · 一切皆插件 · 模型底座之争"),(11,"2026/08/17")]:
        try:
            ph = s.placeholders[pidx]; ph.text = txt
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = dk.EAFONT; dk._apply_ea(r, dk.EAFONT)
        except Exception: pass
    notes(s, "开场:今天用十几分钟,讲清 DeepSeek 开源的 Harness 是什么、为什么重要。")

    # 2 目录(三问骨架)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "目录  Catalog", D.anchor)
    cols = dk.columns(3, slide=s, top=1.4, bottom=0.7, margin=0.5, gap=0.3)
    parts = [("01","厘清产品","WorkBuddy/灵犀/Comate/dsh 各是谁的",["含编程吗 · 含 WPS 办公吗","除了编程还有什么","Comate 内嵌 Pi 的发现"],D.anchor,CARDBG),
             ("02","一切皆插件","dsh 怎么实现一切皆插件",["插座比喻 · 没有特权内核","换模型 = 换一个插件","五层架构 · 锁定底座"],D.comparator,RGBColor(0xFF,0xF1,0xE8)),
             ("03","锁定底座 + 应对","DeepSeek 怎么锁模型底座",["四步打法","我们怎么利用 dsh","为什么爆火 · 是否成标准"],DEEP,RGBColor(0xFF,0xEE,0xF0))]
    for (num,title,q,items,col,fill),c in zip(parts, cols):
        x,y,w,h = c
        card(s, x, y, w, h, fill=fill, line=col, line_w=1.2)
        dk.text(s, x+0.25, y+0.22, w-0.5, 0.4, [[(num, 22, col, True, False, dk.FONT)]], wrap=False)
        dk.text(s, x+0.25, y+0.62, w-0.5, 0.4, [[(title, 16, DARK, True, False, dk.EAFONT)]], wrap=False)
        dk.text(s, x+0.25, y+1.05, w-0.5, 0.7, [[(q, 12.5, D.neutral, False, False, dk.EAFONT)]], line_spacing=1.2)
        yy = y+1.95
        for it in items:
            dk.text(s, x+0.35, yy, w-0.7, 0.32, [[("#  "+it, 12.5, INK, False, False, dk.EAFONT)]]); yy += 0.42
    notes(s, "目录是领导关心的三件事:厘清产品/一切皆插件/锁定底座+应对。讲完答追问:为什么爆火、会不会成标准、我们怎么用。")

    # 3 章节页
    chap(prs, D, "chapter", "01", "先厘清:这些产品到底是什么", sub="WorkBuddy·灵犀·Comate·dsh 各是谁的")
    notes(s, "市面上几个产品常被混为一谈,先把归属讲清楚。")

    # 4 四产品归属表
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "先纠偏:四个产品,四个归属", D.anchor)
    cols = dk.columns(4, slide=s, top=1.35, bottom=1.15, margin=0.5, gap=0.28)
    data = [("DeepSeek Harness","DeepSeek 官方","开源 Agent 框架","让模型干活\n不含办公",D.anchor),
            ("WorkBuddy","腾讯·非DeepSeek","桌面工作台","编程强\n办公靠技能",D.comparator),
            ("WPS 灵犀","金山 WPS","AI 办公智能体","编程+办公\n都在壳里",DEEP),
            ("WPS Comate","金山 WPS","桌面 Agent","内嵌 Pi\n编程+操控文档",D.neutral)]
    for (name,owner,form,feat,col),c in zip(data, cols):
        x,y,w,h = c
        card(s, x, y, w, h, fill=WHITE, line=col, line_w=1.6, r=0.1)
        dk.box(s, x, y, w, 0.52, fill=col, round=True, corners='top', r=0.1)
        dk.text(s, x+0.16, y+0.05, w-0.32, 0.42, [[(name, 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.18, y+0.64, w-0.36, 0.32, [[(owner, 13, col, True, False, dk.EAFONT)]], wrap=False)
        dk.text(s, x+0.18, y+1.0, w-0.36, 0.32, [[(form, 12.5, D.neutral, False, False, dk.EAFONT)]])
        dk.box(s, x+0.18, y+1.38, w-0.36, 0.012, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+0.18, y+1.5, w-0.36, 0.7, [[(feat, 13, INK, False, False, dk.EAFONT)]], line_spacing=1.2)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "关键纠偏", "WorkBuddy 是腾讯的、非 DeepSeek;灵犀(千万办公)才含 WPS 办公功能;Comate 内嵌第三方开源 Pi。", label_c=D.anchor, body_c=DARK)
    notes(s, "WorkBuddy 不是 DeepSeek 的,是腾讯;灵犀才真正含 WPS 办公;Comate 内嵌第三方 Pi。")

    # 4b 除了编程还有什么(问2直答)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "除了编程,WorkBuddy 与灵犀还能干什么?", D.anchor)
    cols = dk.columns(2, slide=s, top=1.35, bottom=1.25, margin=0.5, gap=0.45)
    x,y,w,h = cols[0]
    card(s, x, y, w, h, fill=WHITE, line=D.comparator, line_w=1.5)
    dk.box(s, x, y, w, 0.5, fill=D.comparator, round=True, corners='top', r=0.12)
    dk.text(s, x+0.25, y+0.04, w-0.5, 0.42, [[("WorkBuddy  ·  腾讯", 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, x+0.28, y+0.62, w-0.56, 0.3, [[("编程:",13,D.neutral,True,False,dk.EAFONT),("有(内置 codebuddy 命令行)",13,INK,False,False,dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.28, y+1.0, w-0.56, 0.3, [[("除了编程还有:",13,D.neutral,True,False,dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["多窗口多 Agent 并行","20+ 技能包 + MCP","企微/飞书/钉钉遥控","定时任务 · 140+ 专家角色","金融分析 · 微信支付"]):
        dk.text(s, x+0.4, y+1.34+i*0.34, w-0.7, 0.32, [[("▸ ",11,D.comparator,True,False,dk.FONT),(it,12.5,INK,False,False,dk.EAFONT)]], wrap=False)
    dk.box(s, x+0.28, y+h-0.66, w-0.56, 0.52, fill=RGBColor(0xFF,0xF1,0xE8), line=D.comparator, line_w=1.0, round=True, r=0.08)
    dk.text(s, x+0.4, y+h-0.62, w-0.7, 0.44, [[("WPS 办公:",12,D.comparator,True,False,dk.EAFONT),("不在本体,靠技能间接读写",12.5,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    x,y,w,h = cols[1]
    card(s, x, y, w, h, fill=WHITE, line=DEEP, line_w=1.5)
    dk.box(s, x, y, w, 0.5, fill=DEEP, round=True, corners='top', r=0.12)
    dk.text(s, x+0.25, y+0.04, w-0.5, 0.42, [[("WPS 灵犀  ·  金山(千万办公)", 16, WHITE, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, x+0.28, y+0.62, w-0.56, 0.3, [[("编程:",13,D.neutral,True,False,dk.EAFONT),("有(Python/Node/AirScript)",13,INK,False,False,dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.28, y+1.0, w-0.56, 0.3, [[("除了编程还有:",13,D.neutral,True,False,dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["智能对话(跨会话记忆)","一键生成 Word/PPT/PDF","数据分析 + 可视化图表","网页自动化(内置浏览器)","AI 图像(生图/OCR/音转文)"]):
        dk.text(s, x+0.4, y+1.34+i*0.34, w-0.7, 0.32, [[("▸ ",11,DEEP,True,False,dk.FONT),(it,12.5,INK,False,False,dk.EAFONT)]], wrap=False)
    dk.box(s, x+0.28, y+h-0.66, w-0.56, 0.52, fill=RGBColor(0xFF,0xEE,0xF0), line=DEEP, line_w=1.0, round=True, r=0.08)
    dk.text(s, x+0.4, y+h-0.62, w-0.7, 0.44, [[("WPS 办公:",12,DEEP,True,False,dk.EAFONT),("在,直处理 Word/Excel/PPT/PDF",12.5,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话直答", "含不含 WPS 办公——只有灵犀真正含;WorkBuddy 靠技能间接读写,Comate 经 JSAPI 操控(API 级)。", label_c=D.anchor, body_c=DARK)
    notes(s, "直答问2:除了编程还有什么、WPS功能在不在。WorkBuddy 腾讯,编程有,还有多Agent/技能/遥控/金融,办公不在本体靠技能间接读写。灵犀金山,编程有,还有对话/文档/数据/网页/图像,办公在直处理。只有灵犀真含WPS办公。")

    # 5 关键发现(深色页)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("dark")]); set_title(s, "关键发现:WPS Comate 内嵌的就是 Pi", WHITE)
    cols = dk.columns(2, slide=s, top=1.35, bottom=1.3, margin=0.5, gap=0.5)
    x,y,w,h = cols[0]
    card(s, x, y, w, h, fill=RGBColor(0x2A,0x14,0x18), line=D.anchor, line_w=1.5)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.5, [[("本机实测", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, x+0.3, y+0.6, w-0.6, 1.2, [[("金山自家的 WPS Comate,",16,WHITE,True,False,dk.EAFONT),("核心就是 Pi",16,D.emphasis,True,False,dk.EAFONT),(" v0.79.3",15,WHITE,False,False,dk.EAFONT)]], line_spacing=1.3)
    x,y,w,h = cols[1]
    card(s, x, y, w, h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=1.2)
    dk.text(s, x+0.3, y+0.25, w-0.6, 0.4, [[("意味着什么", 13, D.emphasis, True, False, dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["Pi 轻量路线已被产品化落地","选型与产品方向一致、风险低","与 dsh 同源 pi-ai,可平滑评估"]):
        dk.text(s, x+0.3, y+0.7+i*0.56, w-0.6, 0.5, [[("▸  ",13,D.anchor,True,False,dk.FONT),(it,13,WHITE,False,False,dk.EAFONT)]], line_spacing=1.2)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "三大 harness 互不内嵌", "WorkBuddy 用 codebuddy、Comate 用 Pi、DeepSeek 用 dsh——各自独立。", label_c=D.emphasis, body_c=WHITE, fill=RGBColor(0x2A,0x14,0x18))
    notes(s, "金山自家 Comate 的核心运行时就是第三方 Pi;我们也在用 Pi。三大 harness 互不内嵌。")

    # 6 章节页
    chap(prs, D, "chapter", "02", 'DeepSeek Harness 如何"一切皆插件"', sub="Cordis·能力 seam·插件树·创造模式")
    notes(s, "进入架构,目标让领导看懂插件化体现在哪。")

    # 7 定位+四模式
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "定位:Agent = 模型 + Harness", D.anchor)
    dk.text(s, 0.7, 1.5, 12, 1.2, [[("Agent",40,D.anchor,True,False,dk.FONT),(" = ",32,DARK,True,False,dk.FONT),("模型",34,D.neutral,True,False,dk.EAFONT),(" + ",30,DARK,True,False,dk.FONT),("Harness",40,D.comparator,True,False,dk.FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, 0.7, 2.7, 12, 0.5, [[("模型负责『想』,Harness 负责理解环境、调工具、组织多步任务", 14, MUTE, False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
    modes = [("标准","默认全套工具,日常开发"),("PTC","生成代码组合多步调用"),("极简","只留 Shell,跑基准"),("创造","agent 改自己的插件")]
    cols = dk.columns(4, slide=s, top=3.6, bottom=1.3, margin=0.5, gap=0.3)
    b8 = Build(s)
    for i,((n,d),c) in enumerate(zip(modes, cols)):
        x,y,w,h = c
        with b8.step():
            card(s, x, y, w, h, fill=CARDBG, line=D.neutral, line_w=1.0, r=0.1)
            num_circle(s, x+w/2-0.28, y+0.18, 0.56, i+1, D.anchor, D.comparator, size=16)
            dk.text(s, x+0.12, y+0.82, w-0.24, 0.35, [[(n, 15, D.anchor, True, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
            dk.text(s, x+0.12, y+1.18, w-0.24, 0.6, [[(d, 14, INK, False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    b8.apply(effect="fade")
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "让大模型『能动手干活』的运行时底座,对标 Claude Code / Codex,开源。", label_c=D.anchor, body_c=DARK)
    notes(s, "Agent=模型+Harness。四模式:标准/PTC/极简/创造。创造模式 agent 能改自己的插件。")


    # 7b 一切皆插件(大白话:插座比喻)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "一切皆插件:dsh 是一块插座板,能力都是插头", D.anchor)
    box_x, box_y, box_w, box_h = 0.9, 1.7, 5.0, 4.3
    dk.box(s, box_x, box_y, box_w, box_h, fill=RGBColor(0xF3,0xF5,0xF8), line=D.neutral, line_w=2.5, round=True, r=0.06)
    dk.text(s, box_x+0.2, box_y+0.15, box_w-0.4, 0.4, [[("ctx  ——  服务插座板", 16, D.neutral, True, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
    dk.text(s, box_x+0.2, box_y+0.55, box_w-0.4, 0.3, [[("(按名字找能力,不认具体实现)", 11.5, MUTE, False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER)
    slots = [("ctx.llm","模型","想",D.anchor),("ctx.tools","工具","动手",D.comparator),("ctx.shell","命令行/沙箱","执行",DEEP),("ctx.sessions","记忆/会话","记",D.neutral)]
    sy0 = box_y+1.05; sh = 0.62; gap = 0.22
    for i,(ky,n,v,col) in enumerate(slots):
        yy = sy0 + i*(sh+gap)
        dk.box(s, box_x+0.35, yy, box_w-0.7, sh, fill=WHITE, line=col, line_w=1.4, round=True, r=0.1)
        dk.box(s, box_x+0.35, yy, 0.1, sh, fill=col, round=True, corners='left', r=0.1)
        dk.text(s, box_x+0.55, yy+0.06, 1.7, 0.3, [[(ky, 12.5, col, True, False, dk.FONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, box_x+2.35, yy+0.06, 1.2, 0.3, [[(n, 14, INK, True, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, box_x+3.6, yy+0.06, 0.9, 0.3, [[("负责"+v, 12, MUTE, False, False, dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    rx, rw = 6.3, 6.4
    card(s, rx, 1.7, rw, 2.0, fill=RGBColor(0xFF,0xF4,0xF6), line=D.anchor, line_w=1.4, r=0.08)
    dk.text(s, rx+0.25, 1.82, rw-0.5, 0.4, [[("想换模型?", 16, D.anchor, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, rx+0.25, 2.25, rw-0.5, 1.35, [[("拔下 ctx.llm 这个插头,换一个插上去——", 14, INK, False, False, dk.EAFONT)],[("调用方的代码一行都不用改。", 14, D.anchor, True, False, dk.EAFONT)],[("想换沙箱、想接私有环境?同理,换对应插头。", 13.5, INK, False, False, dk.EAFONT)]], line_spacing=1.3)
    card(s, rx, 3.9, rw, 2.1, fill=WHITE, line=D.neutral, line_w=1.2, r=0.08)
    dk.text(s, rx+0.25, 4.02, rw-0.5, 0.4, [[("没有特权内核", 16, D.neutral, True, False, dk.EAFONT)]], wrap=False)
    dk.text(s, rx+0.25, 4.45, rw-0.5, 1.45, [[("连 agent 循环本身都是一个插头——", 14, INK, False, False, dk.EAFONT)],[("想加新能力,不是改框架,是挂一个新插头。", 14, D.neutral, True, False, dk.EAFONT)],[("插头拔掉,副作用自动撤销(不残留垃圾)。", 13.5, INK, False, False, dk.EAFONT)]], line_spacing=1.3)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "dsh 把每项能力都做成可拔插的插头——模型/工具/命令行/记忆都是,连 agent 循环本身也是。", label_c=D.anchor, body_c=DARK)
    notes(s, "大白话讲一切皆插件。dsh 是插座板 ctx,插着模型/工具/命令行/记忆四个插头。换模型=拔下ctx.llm换一个,代码不用改。没有特权内核,连agent循环都是插头,加能力=挂新插头,拔掉副作用自动撤销。底层叫Cordis。")

    # 8 五层结构(简化,大白话)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "五层结构:每一层都是插件,上一层搭下一层", D.anchor)
    layers = [("apps 组装点","dsh 命令行 · 网页界面 · 给别的程序调",DEEP,RGBColor(0xFF,0xEE,0xF0)),("组合包 / 组装","把基础能力打包,用户可在上面叠加自己的",D.anchor,RGBColor(0xFF,0xF1,0xF3)),("能力插件包","模型/工具/命令行/文件/子代理……几十个,可换",D.anchor,RGBColor(0xFF,0xF4,0xF6)),("核心服务接缝","能力的接口位:模型位/工具位/命令位……",D.neutral,RGBColor(0xF2,0xF5,0xF9)),("Cordis 内核","插座板本体:怎么挂插头/怎么找插头/怎么撤销",D.neutral,RGBColor(0xE9,0xEE,0xF5))]
    lx, lw = 0.7, 11.9; ly0 = 1.5; lh = 0.78; gap = 0.18
    for i,(t,d,col,fill) in enumerate(layers):
        y = ly0 + i*(lh+gap)
        card(s, lx, y, lw, lh, fill=fill, line=col, line_w=1.6 if i>=2 else 1.0, r=0.08)
        dk.box(s, lx, y, 0.1, lh, fill=col, round=True, corners='left', r=0.08)
        dk.text(s, lx+0.25, y+0.05, 3.0, 0.34, [[(t,14,col,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.text(s, lx+3.3, y+0.05, lw-3.5, lh-0.1, [[(d,13,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        if i < len(layers)-1: dk.arrow(s, lx+lw/2-0.12, y+lh+0.0, 0.24, gap, color=RGBColor(0xB0,0xB8,0xC4), direction='down')
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "同一内核多种形态", "同一套内核,装不同插件 = 命令行/网页/给别的程序调——还能对外暴露接口被集成。", label_c=D.anchor, body_c=DARK)
    notes(s, "dsh 五层从底到顶:Cordis内核(插座板本体)→核心服务接缝(接口位)→能力插件包(几十个可换)→组合包/组装→apps组装点。组合↑依赖↓。同一内核多种形态。备问:agent循环怎么跑见md 2.7。")

    # 9 换模型=换一个插头(signature move)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "换模型 = 换一个插头:这就是一切皆插件的精髓", D.anchor)
    slot_x, slot_y, slot_w, slot_h = 0.62, 1.4, 3.3, 2.6
    card(s, slot_x, slot_y, slot_w, slot_h, fill=RGBColor(0x1A,0x22,0x2E), line=D.neutral, line_w=2.0, r=0.1)
    dk.text(s, slot_x+0.2, slot_y+0.15, slot_w-0.4, 0.4, [[("ctx.llm 插槽", 14, D.emphasis, True, False, dk.FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    dk.text(s, slot_x+0.2, slot_y+0.55, slot_w-0.4, 0.5, [[("(模型这个能力的接口位)", 11.5, RGBColor(0xBB,0xBB,0xCC), False, False, dk.EAFONT)]])
    dk.box(s, slot_x+0.35, slot_y+1.1, slot_w-0.7, 1.2, fill=RGBColor(0x0D,0x12,0x1A), line=D.anchor, line_w=1.5, round=True, r=0.08)
    dk.text(s, slot_x+0.35, slot_y+1.1, slot_w-0.7, 1.2, [[("模型插头\n插这里", 14, RGBColor(0xFF,0xAA,0xAA), False, False, dk.EAFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    dk.arrow(s, slot_x+slot_w+0.05, slot_y+slot_h/2-0.25, 0.7, 0.5, color=D.anchor, direction='right')
    provs = [("DeepSeek 原生","dsh 默认带它",D.anchor,True),("经 pi-ai 接任意","能接公司 wpsyun 模型",D.comparator,False),("回放测试","重放对话做测试",MUTE,False)]
    px = slot_x+slot_w+0.95; pw = (W_IN-0.62-px-0.3)/3 - 0.2
    for i,(nm,d,col,hi) in enumerate(provs):
        x = px + i*(pw+0.25); y = slot_y+0.1
        card(s, x, y, pw, 2.6, fill=RGBColor(0xFF,0xF4,0xF6) if hi else WHITE, line=col, line_w=2.2 if hi else 1.2, r=0.1)
        if hi:
            dk.box(s, x, y, pw, 0.4, fill=col, round=True, corners='top', r=0.1)
            dk.text(s, x+0.1, y+0.02, pw-0.2, 0.36, [[("★ 默认",11,WHITE,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dk.text(s, x+0.12, y+0.5, pw-0.24, 0.4, [[(nm,14,col,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, wrap=False)
        dk.text(s, x+0.12, y+0.95, pw-0.24, 0.6, [[(d,12.5,INK,False,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.15)
        dk.text(s, x+0.12, y+1.62, pw-0.24, 0.8, [[("换这个\n= 换模型底座",14,D.anchor,True,False,dk.EAFONT)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "模型在 dsh 里就是一个可换的插头——换模型底座,调用方代码一行不改;同理换沙箱、换子代理也是换插头。", label_c=D.anchor, body_c=DARK)
    notes(s, "ctx.llm 是模型接口位,三张卡是可换模型插头:DeepSeek原生(默认)、经pi-ai接任意(能接公司wpsyun)、回放测试。换模型=换一个插头,代码不改。底层叫seam,有三角色,深问见md 2.2。")


    # (extensions 生态页已下沉讲者备注;完整内容见同名 md 第 2.8 节作备问弹药)

    # 10 章节页
    chap(prs, D, "chapter", "03", "锁定模型底座 + 我们的应对", sub="DeepSeek 怎么锁 · 我们怎么用 · 会不会成标准")
    notes(s, "进入战略,今天最该记住的部分。")

    # 11 战略四步
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "战略:用『动手层』标准化,反向锁定『模型层』", D.anchor)
    steps = [("默认带 DeepSeek","模型在 dsh 是可换插件,但默认体验/生态偏 DeepSeek"),("先选 harness 再选模型","团队建技术栈后迁移成本陡增;换模型只是一个插件"),("数据/反馈回流","更多 Agent 跑在 DeepSeek 上,真实数据反哺迭代"),("生态飞轮","Cordis + dsh-plugin 话题;24h 已 288 个插件仓")]
    rows = dk.rows(4, slide=s, top=1.4, bottom=1.25, gap=0.22)
    b13 = Build(s)
    for i,((n,d),c) in enumerate(zip(steps, rows)):
        x,y,w,h = c
        with b13.step():
            num_circle(s, x+0.1, y+h/2-0.27, 0.54, i+1, D.anchor, D.comparator, size=18)
            card(s, x+0.8, y, w-0.8, h, fill=CARDBG, line=RGBColor(0xE8,0xE8,0xEE), line_w=1.0, r=0.08)
            dk.text(s, x+1.0, y+0.12, w-1.2, 0.35, [[(n,15,D.anchor,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            dk.text(s, x+1.0, y+0.12, w-1.2, h-0.24, [[("　",1,WHITE,False,False),(d,14,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    b13.apply(effect="fade")
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "一句话", "谁定义 harness 标准,谁就定义模型被怎么用、用谁的——DeepSeek 同时开源模型与 harness,用动手层标准化反向锁定模型层。", label_c=D.anchor, body_c=DARK)
    notes(s, "四步:默认带DeepSeek/先选harness/数据回流/生态飞轮(早期约288仓,引用以GitHub为准)。谁定义harness谁定义模型分发。")

    # 12 预判追问(爆火/标准)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "讲完你可能追问:为什么爆火 · 会不会成中国标准", D.anchor)
    dk.text(s, 0.7, 1.32, 12, 0.34, [[("追问 ①  为什么 dsh 爆火,Pi 却没有?", 14.5, D.anchor, True, False, dk.EAFONT)]], wrap=False)
    cols = dk.columns(2, slide=s, top=1.75, bottom=3.55, margin=0.7, gap=0.4)
    x,y,w,h = cols[0]
    card(s, x, y, w, h, fill=RGBColor(0xFF,0xF4,0xF6), line=D.anchor, line_w=1.2, r=0.08)
    dk.text(s, x+0.2, y+0.12, w-0.4, 0.3, [[("dsh 爆火", 13, D.anchor, True, False, dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["DeepSeek 品牌势能(V 系列已火过一次)","官方全开源 + 有网页界面,门槛低","官方维护几十个能力包,开箱即用","架构叙事高(Cordis + 论文),有记忆点"]):
        dk.text(s, x+0.3, y+0.5+i*0.42, w-0.55, 0.4, [[("▸ ",11,D.anchor,True,False,dk.FONT),(it,12,INK,False,False,dk.EAFONT)]], line_spacing=1.15)
    x,y,w,h = cols[1]
    card(s, x, y, w, h, fill=WHITE, line=D.comparator, line_w=1.2, r=0.08)
    dk.text(s, x+0.2, y+0.12, w-0.4, 0.3, [[("Pi 没同样爆火", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    pi_pts = [("独立开发者个人项目,无品牌势能",INK,False),("纯终端、无网页界面,门槛高",INK,False),("极简哲学,把很多能力推给生态",INK,False),("但:走被集成路线,金山 Comate 已内嵌 Pi",D.anchor,True)]
    for i,(it,col,bd) in enumerate(pi_pts):
        dk.text(s, x+0.3, y+0.5+i*0.42, w-0.55, 0.4, [[("▸ ",11,D.comparator,True,False,dk.FONT),(it,12,col,bd,False,dk.EAFONT)]], line_spacing=1.15)
    dk.text(s, 0.7, 4.05, 12, 0.34, [[("追问 ②  dsh 会成为 AI Agent 的中国标准吗?", 14.5, D.anchor, True, False, dk.EAFONT)]], wrap=False)
    cols2 = dk.columns(2, slide=s, top=4.5, bottom=1.25, margin=0.7, gap=0.4)
    x,y,w,h = cols2[0]
    card(s, x, y, w, h, fill=RGBColor(0xEE,0xF5,0xEE), line=RGBColor(0x2E,0x8B,0x57), line_w=1.2, r=0.08)
    dk.text(s, x+0.2, y+0.12, w-0.4, 0.3, [[("支持:最强候选", 13, RGBColor(0x2E,0x8B,0x57), True, False, dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["唯一中国头部 AI 公司官方全开源 harness","seam 化适配政企换模型/换部署/换沙箱","多 provider 不硬绑,利于成公共标准"]):
        dk.text(s, x+0.3, y+0.5+i*0.4, w-0.55, 0.38, [[("▸ ",11,RGBColor(0x2E,0x8B,0x57),True,False,dk.FONT),(it,12,INK,False,False,dk.EAFONT)]], line_spacing=1.15)
    x,y,w,h = cols2[1]
    card(s, x, y, w, h, fill=RGBColor(0xFF,0xF1,0xE8), line=D.comparator, line_w=1.2, r=0.08)
    dk.text(s, x+0.2, y+0.12, w-0.4, 0.3, [[("存疑:现在说会太早", 13, D.comparator, True, False, dk.EAFONT)]], wrap=False)
    for i,it in enumerate(["开发者预览,有破坏性变更,未稳定","大厂各有 agent 框架,标准之争持续","标准要靠企业落地,不是 star 数"]):
        dk.text(s, x+0.3, y+0.5+i*0.4, w-0.55, 0.38, [[("▸ ",11,D.comparator,True,False,dk.FONT),(it,12,INK,False,False,dk.EAFONT)]], line_spacing=1.15)
    bottom_callout_at(s, 0.5, W_IN-1.0, 7.07, "结论", "当前最有力的中国开源 agent 事实标准候选——但是否真成标准取决于 1-2 年内企业落地与版本稳定,现在未定。", label_c=D.anchor, body_c=DARK)
    notes(s, "追问①:dsh 爆火靠品牌势能/全开源有界面/官方能力包/架构叙事;Pi 独立开发者无品牌、纯终端、极简哲学,但走被集成路线(Comate内嵌)。star≠标准。追问②:最强候选(开源/seam化/不硬绑)但未定(预览/大厂竞争/需企业落地)。完整见md第7节。")

    # 13 我们怎么利用 dsh(问3直答)
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "我们怎么利用 dsh:五条应对(能做什么)", D.anchor)
    advs = [("借鉴 seam 化","把模型/文档/沙箱做成可换插件 → 适配私有化换模型换部署换沙箱","借鉴非采用"),("接 wpsyun 跑 PoC","用 dsh + pi-ai 接公司模型,验证在第三方框架可用","短期可做"),("轻 + 重 组合","Pi 做单兵/Comate 级产品,dsh 作平台化底座备选","已验证+观察"),("跟踪标准 · 不上生产","跟踪避免被动;开发者预览,定位预研/POC","版本收敛再评估"),("厘清口径","对外统一归属口径,讲混损害方案可信度","沟通纪律")]
    rows = dk.rows(5, slide=s, top=1.4, bottom=0.75, gap=0.2)
    for i,((n,d,how),c) in enumerate(zip(advs, rows)):
        x,y,w,h = c
        num_circle(s, x, y+h/2-0.24, 0.48, i+1, D.anchor, D.comparator, size=15)
        card(s, x+0.65, y, w-0.65, h, fill=WHITE if i%2==0 else CARDBG, line=RGBColor(0xE8,0xE8,0xEE), line_w=1.0, r=0.06)
        dk.text(s, x+0.85, y+0.05, 3.1, h-0.1, [[(n,14,D.anchor,True,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        dk.box(s, x+4.05, y+0.1, 0.012, h-0.2, fill=RGBColor(0xDD,0xDD,0xDD))
        dk.text(s, x+4.25, y+0.05, w-4.45, h-0.1, [[(d,14,INK,False,False,dk.EAFONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    dk.text(s, 0.62, H_IN-0.42, 12, 0.32, [[("口径:",11,MUTE,False,False,dk.EAFONT),("WorkBuddy=腾讯·灵犀=金山·Comate=金山(内嵌Pi)·Pi=第三方·dsh=DeepSeek",11,D.neutral,False,False,dk.EAFONT)]])
    notes(s, "直答我们怎么利用dsh:借鉴seam化/接wpsyun跑PoC/轻+重组合/跟踪标准不上生产/厘清口径。短期Pi坚定,dsh是中期变量。")

    # 14 结论
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("red_conclusion")]); set_title(s, "结论", WHITE)
    dk.text(s, 1.0, 1.9, 11.3, 2.8, [[("dsh 既是",24,WHITE,True,False,dk.EAFONT),("技术借鉴对象",26,D.emphasis,True,False,dk.EAFONT),(",也是",24,WHITE,True,False,dk.EAFONT),("战略变量",26,D.emphasis,True,False,dk.EAFONT),("。",24,WHITE,True,False,dk.EAFONT)],[("",8,WHITE,False,False)],[("短期:Pi 路线已被验证落地、且本团队在用——可坚定。",16,WHITE,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("中期:dsh 作平台化底座候选持续观察,借鉴其 seam 化设计。",16,WHITE,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("长期:评估把自身方案架构『接缝化』,对冲模型层趋同风险。",16,WHITE,False,False,dk.EAFONT)],[("",10,WHITE,False,False)],[("三追问(爆火/中国标准/我们怎么用)速答见同名 md 第 7 节。",12.5,D.emphasis,True,False,dk.EAFONT)]], line_spacing=1.2)
    notes(s, "收尾:dsh 既是借鉴对象也是战略变量。短期Pi坚定,中期dsh观察,长期接缝化。")

    # 15 附录
    s = prs.slides.add_slide(prs.slide_layouts[D.P.layout("content")]); set_title(s, "附录·证据出处", D.anchor)
    dk.text(s, 0.7, 1.5, 12, 5, [[("A. DeepSeek Harness(本地源码):",12.5,D.anchor,True,False,dk.EAFONT),(" README/architecture/cordis-primer · agent-presets · cordis-host-runner · llm-pi-ai",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("B. Pi(本机已装 0.84.2):",12.5,D.comparator,True,False,dk.EAFONT),(" @earendil-works/pi-coding-agent · 本机 models.json · 与 dsh 共用 pi-ai",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("C. WorkBuddy(本机实测):",12.5,D.anchor,True,False,dk.EAFONT),(" AppData/Local/Programs/WorkBuddy(腾讯,Electron)·codebuddy CLI·copilot.tencent.com",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("D. WPS Comate(本机实测):",12.5,D.anchor,True,False,dk.EAFONT),(" AppData/Local/WPS Comate(金山·内嵌Pi v0.79.3)·comate.wps.cn",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("E. WPS 灵犀(官方):",12.5,D.anchor,True,False,dk.EAFONT),(" lingxi.wps.cn · Python/Node/AirScript + WPS Office 集成",11,INK,False,False,dk.EAFONT)],[("",6,WHITE,False,False)],[("完整证据见同名 md 培训稿附录。",11,MUTE,False,False,dk.EAFONT)]], line_spacing=1.1)
    notes(s, "附录:所有结论的证据出处。完整证据链在培训稿 md 附录。")

    dk.lint_layout(prs, strict=True)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
